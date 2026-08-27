#!/usr/bin/env python3
"""Build the talk deck from talk_content.py.

Assembly only: not a word of the narrative lives here. The design is deliberately plain, because
the figures already carry three colours with meaning (blue both cultivars, orange only the mutant,
red neither) and a template that adds a fourth would compete with them.

Type scale and layout follow the assertion-evidence checklist: a full-sentence title of at most two
lines, evidence rather than bullets underneath, and nothing on the slide the speaker is going to
read out loud. Scenario colours are the IPCC AR6 ones so anyone who has seen an AR6 figure
recognises them without a legend.

Every number reaching a slide comes from one of the six metric tables load_numbers() reads. Almost
none is typed into the content file, and the exceptions are marked in its own docstring, so a slide
cannot drift from the table that produced it.

Usage: python 35_build_talk_pptx.py [--out ../03_presentacion/charla_plinius.pptx]
Requires: python-pptx, Pillow. Run 31 to 34 and 36 to 43 first: the figures, the GIF and the metric
tables must all exist.
"""

import argparse
import csv
import re
import sys
from pathlib import Path

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
import talk_content  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent

# --- design system ---------------------------------------------------------------------------
W, H = Inches(13.333), Inches(7.5)          # 16:9
MARGIN = Inches(0.62)
CONTENT_W = W - 2 * MARGIN

INK = RGBColor(0x16, 0x18, 0x1D)
MUTED = RGBColor(0x5F, 0x66, 0x72)
FAINT = RGBColor(0x8A, 0x90, 0x99)
RULE = RGBColor(0xDC, 0xDF, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BLUE = RGBColor(0x2C, 0x7B, 0xB6)           # both cultivars viable
ORANGE = RGBColor(0xFD, 0xAE, 0x61)         # only the mutant
RED = RGBColor(0xD7, 0x19, 0x1C)            # neither

FONT = "Segoe UI"
TITLE_PT, SUB_PT, BODY_PT, CAP_PT, SRC_PT = 26, 15, 16, 13, 10


def txbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    """A text box with autofit off and no internal padding, which is never what a slide wants."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def write(tf, text, *, size, colour=INK, bold=False, para=0, space_after=0, align=PP_ALIGN.LEFT,
          line=None):
    p = tf.paragraphs[para] if para < len(tf.paragraphs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = colour
    return p


def link_urls(p, *, size, colour=MUTED):
    """Rebuild a finished paragraph so that any https:// token becomes a real hyperlink.

    Splitting the text into runs is what makes the address clickable in the exported PDF, which is
    how the repository actually gets visited after a talk. The colour is deliberately left as the
    surrounding text: a blue underlined link would be the only one in the whole deck.
    """
    text = "".join(r.text for r in p.runs)
    if "https://" not in text:
        return
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    for token in re.split(r"(https://\S+)", text):
        if not token:
            continue
        run = p.add_run()
        run.text = token
        run.font.size = Pt(size)
        run.font.name = FONT
        run.font.color.rgb = colour
        if token.startswith("https://"):
            run.hyperlink.address = token


def para(tf, text, **kw):
    """Append a paragraph to an existing frame."""
    tf.add_paragraph()
    return write(tf, text, para=len(tf.paragraphs) - 1, **kw)


def rect(slide, x, y, w, h, fill, *, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    return sh


def fit(img_path, box_w, box_h):
    """Largest (w, h) fitting the box while preserving the image's aspect ratio."""
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    return int(iw * scale), int(ih * scale)


def place_image(slide, path, x, y, box_w, box_h):
    """Centre an image inside a box, scaled to fit. Returns its bounding box."""
    if not path.exists():
        sys.exit(f"missing figure: {path}\n  run scripts 31-34 before building the deck")
    w, h = fit(path, box_w, box_h)
    left = x + int((box_w - w) / 2)
    top = y + int((box_h - h) / 2)
    slide.shapes.add_picture(str(path), left, top, width=w, height=h)
    return left, top, w, h


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def fade(slide, ms=350):
    """Give the slide a short fade-in, written straight into the slide XML.

    python-pptx has no transition API, and applying transitions by hand in PowerPoint afterwards
    would mean losing them every time this script regenerates the file. Writing the element here
    keeps the deck reproducible from source, which is the whole point of building it with a script.

    Schema order inside <p:sld> is cSld, clrMapOvr, transition, timing, so the element is appended
    after clrMapOvr rather than at the end.
    """
    from lxml import etree
    sld = slide._element
    tr = etree.SubElement(sld, f"{{{P_NS}}}transition")
    tr.set("spd", "fast")
    tr.set("advClick", "1")
    etree.SubElement(tr, f"{{{P_NS}}}fade")
    timing = sld.find(f"{{{P_NS}}}timing")
    if timing is not None:                      # keep timing last if python-pptx already made one
        sld.remove(timing)
        sld.append(timing)
    return tr


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fade(slide)
    return slide


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def title_block(slide, text, y=None):
    """Assertion title plus the hairline under it. Returns the y where content may start."""
    y = MARGIN if y is None else y
    lines = 2 if len(text) > 62 else 1
    h = Inches(0.52 * lines)
    tf = txbox(slide, MARGIN, y, CONTENT_W, h)
    write(tf, text, size=TITLE_PT, bold=True, line=0.92)
    ry = y + h + Inches(0.12)
    rect(slide, MARGIN, ry, CONTENT_W, Emu(9525), RULE)
    return ry + Inches(0.22)


def footer(slide, source, page):
    if source:
        tf = txbox(slide, MARGIN, H - Inches(0.5), CONTENT_W - Inches(0.6), Inches(0.3))
        write(tf, source, size=SRC_PT, colour=FAINT)
    tf = txbox(slide, W - MARGIN - Inches(0.6), H - Inches(0.5), Inches(0.6), Inches(0.3))
    write(tf, str(page), size=SRC_PT, colour=FAINT, align=PP_ALIGN.RIGHT)


# --- slide kinds -----------------------------------------------------------------------------
def place_logos(slide, paths, y, *, h=Inches(0.42), gap=Inches(0.34), x=None):
    """A row of institutional logos, each scaled to a common height rather than a common width.

    Logos arrive at whatever aspect their owner published, so matching widths would make a square
    mark tower over a horizontal lockup. Matching the cap height is what makes a row of them look
    deliberate. Returns the x the row ended at.
    """
    x = MARGIN if x is None else x
    for rel in paths:
        path = ROOT / rel
        if not path.exists():
            sys.exit(f"missing logo: {path}\n  see assets/README.md")
        with Image.open(path) as im:
            iw, ih = im.size
        w = int(h * iw / ih)
        slide.shapes.add_picture(str(path), x, y, width=w, height=h)
        x += w + gap
    return x


def s_cover(prs, d, page):
    # No accent bar and a 32 pt title: the working title is a full descriptive sentence rather
    # than a short phrase, so it needs three lines and the whole upper band of the page.
    slide = blank(prs)
    tf = txbox(slide, MARGIN, Inches(0.68), CONTENT_W - Inches(1.2), Inches(2.05))
    write(tf, d["title"], size=32, bold=True, line=0.95)
    tf = txbox(slide, MARGIN, Inches(3.25), CONTENT_W - Inches(2.0), Inches(1.1))
    write(tf, d["subtitle"], size=17, colour=MUTED, line=1.15)
    tf = txbox(slide, MARGIN, Inches(4.85), CONTENT_W, Inches(1.3))
    write(tf, d["authors"], size=16, bold=True)
    # The affiliations arrive as a list when there are more than two, because running three
    # institutions onto one line is what made the earlier cover unreadable at the back of a room.
    for line in (d["affil"] if isinstance(d["affil"], (list, tuple)) else [d["affil"]]):
        para(tf, line, size=11.5, colour=MUTED, space_after=0)
    rect(slide, MARGIN, Inches(6.15), Inches(1.4), Emu(19050), BLUE)
    tf = txbox(slide, MARGIN, Inches(6.42), CONTENT_W, Inches(0.4))
    write(tf, d["venue"], size=12, colour=MUTED)
    if d.get("logos"):
        place_logos(slide, d["logos"], Inches(6.36), h=Inches(0.46),
                    x=W - MARGIN - Inches(4.4))
    notes(slide, d["notes"])


def s_section(prs, d, page):
    """A divider carrying where the talk has got to and what the section answers.

    A divider that only announces a heading spends a whole page on four words. `parts` and `at`
    draw the run of sections with the current one filled, so a reader who put the deck down knows
    where they are, and `asks` states the questions the section is about to answer, which is what
    a listener actually needs before the evidence starts.
    """
    slide = blank(prs)
    rect(slide, Emu(0), Emu(0), W, H, RGBColor(0xF4, 0xF6, 0xF8))
    rect(slide, Emu(0), Emu(0), Inches(0.22), H, BLUE)
    tf = txbox(slide, Inches(1.1), Inches(2.06), Inches(1.2), Inches(1.4))
    write(tf, d["n"], size=76, bold=True, colour=RGBColor(0xC9, 0xD3, 0xDC))
    tf = txbox(slide, Inches(2.35), Inches(2.18), Inches(9.4), Inches(1.0))
    write(tf, d["title"], size=34, bold=True)
    y = Inches(3.28)
    lead = d.get("lead", "")
    if lead:
        h = est_h(lead, Inches(8.6), 16, line=1.2)
        tf = txbox(slide, Inches(2.35), y, Inches(8.6), h)
        write(tf, lead, size=16, colour=MUTED, line=1.2)
        y += h + Inches(0.22)
    for ask in d.get("asks", ()):
        rect(slide, Inches(2.35), y + Inches(0.07), Inches(0.09), Inches(0.09), BLUE)
        h = est_h(ask, Inches(8.2), 13.5, line=1.2)
        tf = txbox(slide, Inches(2.60), y, Inches(8.2), h)
        write(tf, ask, size=13.5, colour=INK, line=1.2)
        y += h + Inches(0.13)

    # The run of sections along the foot: filled to here, hollow after it. Drawn from the same list
    # on every divider, so a part that gets renamed cannot disagree with itself two slides later.
    parts, at = d.get("parts"), d.get("at")
    if parts and at:
        pill_y = H - Inches(0.92)
        gap = Inches(0.10)
        pw = int((CONTENT_W - gap * (len(parts) - 1)) / len(parts))
        for i, name in enumerate(parts, start=1):
            px = MARGIN + (i - 1) * (pw + gap)
            done = i <= at
            rect(slide, px, pill_y, pw, Inches(0.055), BLUE if done else RULE)
            tf = txbox(slide, px, pill_y + Inches(0.14), pw, Inches(0.3))
            write(tf, name, size=9, bold=(i == at), colour=INK if i == at else FAINT)
    notes(slide, d.get("notes", ""))


def s_figure(prs, d, page):
    slide = blank(prs)
    y = title_block(slide, d["title"])
    cap = d.get("caption")
    cap_h = Inches(0.62) if cap else Inches(0)
    box_h = H - y - Inches(0.72) - cap_h
    left, top, w, h = place_image(slide, ROOT / d["image"], MARGIN, y, CONTENT_W, box_h)
    if cap:
        tf = txbox(slide, MARGIN, y + box_h + Inches(0.06), CONTENT_W, cap_h)
        write(tf, cap, size=CAP_PT, colour=MUTED, line=1.15)
    if d.get("gif"):
        # PowerPoint only animates a GIF in slideshow mode; in the editing view it shows frame one.
        badge = rect(slide, left + w - Inches(1.65), top + Inches(0.08), Inches(1.55),
                     Inches(0.28), RGBColor(0x16, 0x18, 0x1D))
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(tf, "animated · press F5", size=9, colour=WHITE, align=PP_ALIGN.CENTER)
    footer(slide, d.get("source"), page)
    notes(slide, d["notes"])


def s_figure_side(prs, d, page):
    slide = blank(prs)
    y = title_block(slide, d["title"])
    box_h = H - y - Inches(0.72)
    img_w = int(CONTENT_W * 0.56)
    place_image(slide, ROOT / d["image"], MARGIN, y, img_w, box_h)
    tx = MARGIN + img_w + Inches(0.45)
    tw = W - MARGIN - tx
    tf = txbox(slide, tx, y + Inches(0.1), tw, box_h, anchor=MSO_ANCHOR.TOP)
    for i, pt in enumerate(d["points"]):
        write(tf, pt, size=BODY_PT - 2, colour=INK if i == 0 else MUTED, para=i,
              space_after=13, line=1.18) if i == 0 else para(
                  tf, pt, size=BODY_PT - 2, colour=MUTED, space_after=13, line=1.18)
    footer(slide, d.get("source"), page)
    notes(slide, d["notes"])


def s_compare(prs, d, page):
    slide = blank(prs)
    y = title_block(slide, d["title"])
    col_w = int((CONTENT_W - Inches(0.5)) / 2)
    # The card is sized from its longest column rather than fixed. A fixed 3.55 in held three
    # short lines and clipped the fourth as soon as one of them wrapped.
    n_lines = max(len(d["left"]["lines"]), len(d["right"]["lines"]))
    card_h = Inches(2.15 + 0.42 * n_lines)
    for i, (side, accent) in enumerate([(d["left"], RED), (d["right"], ORANGE)]):
        x = MARGIN + i * (col_w + Inches(0.5))
        card = rect(slide, x, y, col_w, card_h, RGBColor(0xF7, 0xF9, 0xFB))
        card.text_frame.text = ""
        rect(slide, x, y, col_w, Inches(0.09), accent)
        tf = txbox(slide, x + Inches(0.38), y + Inches(0.4), col_w - Inches(0.76), Inches(0.5))
        write(tf, side["head"], size=21, bold=True)
        tf = txbox(slide, x + Inches(0.38), y + Inches(0.94), col_w - Inches(0.76), Inches(0.85))
        write(tf, side["big"], size=46, bold=True, colour=accent)
        tf = txbox(slide, x + Inches(0.38), y + Inches(1.95),
                   col_w - Inches(0.76), card_h - Inches(2.1))
        for j, ln in enumerate(side["lines"]):
            (write if j == 0 else para)(tf, ln, size=12.5, colour=MUTED, space_after=6, line=1.12,
                                        **({"para": 0} if j == 0 else {}))
    tf = txbox(slide, MARGIN, y + card_h + Inches(0.24), CONTENT_W, Inches(0.8))
    write(tf, d["foot"], size=BODY_PT, colour=INK, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d["notes"])


def s_ingredients(prs, d, page):
    slide = blank(prs)
    y = title_block(slide, d["title"])
    n = len(d["items"])
    gap = Inches(0.36)
    col_w = int((CONTENT_W - gap * (n - 1)) / n)
    accents = [BLUE, ORANGE, RGBColor(0x4D, 0x9A, 0x53), RGBColor(0x7B, 0x5E, 0xA7)]
    for i, item in enumerate(d["items"]):
        x = MARGIN + i * (col_w + gap)
        rect(slide, x, y, col_w, Inches(3.3), RGBColor(0xF7, 0xF9, 0xFB))
        rect(slide, x, y, col_w, Inches(0.09), accents[i % len(accents)])
        tf = txbox(slide, x + Inches(0.32), y + Inches(0.44), col_w - Inches(0.64), Inches(0.5))
        write(tf, item["head"], size=20, bold=True)
        tf = txbox(slide, x + Inches(0.32), y + Inches(1.12), col_w - Inches(0.64), Inches(2.0))
        for j, ln in enumerate(item["body"].split("\n")):
            (write if j == 0 else para)(tf, ln, size=13.5, colour=MUTED, space_after=8, line=1.15,
                                        **({"para": 0} if j == 0 else {}))
    tf = txbox(slide, MARGIN, y + Inches(3.7), CONTENT_W, Inches(0.8))
    write(tf, d["foot"], size=BODY_PT, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d["notes"])


def s_close(prs, d, page):
    slide = blank(prs)
    rect(slide, Emu(0), Emu(0), W, Inches(0.14), BLUE)
    tf = txbox(slide, MARGIN, Inches(0.85), CONTENT_W, Inches(0.9))
    write(tf, d["title"], size=38, bold=True)
    # Each point used to be given a fixed 1.06 in whatever it said, so a fifth point of two and a
    # half lines ran under the repository line at the foot. The points are measured now and the
    # type steps down until the block fits the band between the title and that line.
    y = Inches(2.15)
    w = CONTENT_W - Inches(0.34)
    avail = H - Inches(1.10) - y
    pts, gap, size = d["points"], Inches(0.30), 17
    while size > 12:
        heights = [est_h(p, w, size, line=1.18, pad=0.06) for p in pts]
        if sum(heights) + gap * (len(pts) - 1) <= avail:
            break
        size -= 0.5
    for pt, h in zip(pts, heights):
        rect(slide, MARGIN, y + Inches(0.06), Inches(0.075), min(h, Inches(0.5)), BLUE)
        tf = txbox(slide, MARGIN + Inches(0.34), y, w, h)
        write(tf, pt, size=size, line=1.18)
        y += h + gap
    tf = txbox(slide, MARGIN, H - Inches(0.95), CONTENT_W, Inches(0.4))
    p = write(tf, d["foot"], size=12.5, colour=MUTED)
    link_urls(p, size=12.5)
    notes(slide, d["notes"])


def cover_crop(path, box_w, box_h):
    """A copy of `path` centre-cropped to the aspect of the box, so it fills without stretching.

    place_image() fits an image inside its box, which is right for a figure and wrong for a
    photograph used as a panel: a 3:2 photo fitted into a tall panel leaves two thirds of it white.
    Cropping to the box's aspect first means the panel is filled and nothing is distorted, at the
    cost of losing the edges of the frame, which for a background photograph is what you want.
    """
    want = box_w / box_h
    out = TRIM_DIR / f"{path.stem}_fill{round(want * 1000)}{path.suffix}"
    if out.exists() and out.stat().st_mtime >= path.stat().st_mtime:
        return out
    TRIM_DIR.mkdir(parents=True, exist_ok=True)
    with Image.open(path) as im:
        im = im.convert("RGB")
        w, h = im.size
        if w / h > want:                       # too wide: trim the sides
            new_w = int(h * want)
            box = ((w - new_w) // 2, 0, (w - new_w) // 2 + new_w, h)
        else:                                  # too tall: trim top and bottom
            new_h = int(w / want)
            box = (0, (h - new_h) // 2, w, (h - new_h) // 2 + new_h)
        im.crop(box).save(out, quality=92)
    return out


def s_closing(prs, d, page):
    """The last page of a talk: thanks and contact on the left, the crop on the right.

    Deliberately the only slide in the deck with a photograph. Everything before it argues from
    figures, and a room that has spent fifteen minutes reading maps should be looking at the tree
    the maps are about when the questions start.
    """
    slide = blank(prs)
    photo_w = Inches(5.35)
    photo_x = W - photo_w
    if d.get("photo"):
        src = ROOT / d["photo"]
        if not src.exists():
            sys.exit(f"missing closing photo: {src}\n  see assets/README.md")
        slide.shapes.add_picture(str(cover_crop(src, photo_w, H)), photo_x, Emu(0),
                                 width=photo_w, height=H)
    rect(slide, photo_x - Inches(0.06), Emu(0), Inches(0.06), H, BLUE)

    left_w = photo_x - MARGIN - Inches(0.55)
    rect(slide, MARGIN, Inches(1.30), Inches(1.4), Emu(19050), BLUE)
    tf = txbox(slide, MARGIN, Inches(1.62), left_w, Inches(0.9))
    write(tf, d["title"], size=40, bold=True, line=0.98)

    y = Inches(2.72)
    for pt in d.get("points", ()):
        h = est_h(pt, left_w - Inches(0.26), 13, line=1.22)
        rect(slide, MARGIN, y + Inches(0.06), Inches(0.09), Inches(0.09), BLUE)
        tf = txbox(slide, MARGIN + Inches(0.26), y, left_w - Inches(0.26), h)
        write(tf, pt, size=13, colour=INK, line=1.22)
        y += h + Inches(0.20)

    # The contact line is the one thing on this slide somebody might write down, so it sits on a
    # tinted plate rather than loose in the margin, and the address is a live link in the PDF.
    foot = d.get("foot")
    if foot:
        fh = est_h(foot, left_w - Inches(0.36), 12.5, line=1.24) + Inches(0.26)
        rrect(slide, MARGIN, y + Inches(0.16), left_w, fh, PANEL)
        tf = txbox(slide, MARGIN + Inches(0.18), y + Inches(0.29), left_w - Inches(0.36),
                   fh - Inches(0.26))
        p = write(tf, foot, size=12.5, colour=INK, line=1.24)
        link_urls(p, size=12.5, colour=INK)

    if d.get("logos"):
        place_logos(slide, d["logos"], H - Inches(1.05), h=Inches(0.5))
    notes(slide, d.get("notes", ""))


def s_gallery(prs, d, page):
    """Contact sheet: every map at thumbnail size, so one can be pointed at during questions."""
    slide = blank(prs)
    y = title_block(slide, d["title"])
    items = d["items"]
    cols, rows = 5, 3
    gap = Inches(0.14)
    cell_w = int((CONTENT_W - gap * (cols - 1)) / cols)
    avail_h = H - y - Inches(0.72)
    cell_h = int((avail_h - gap * (rows - 1)) / rows)
    lab_h = Inches(0.23)
    for i, (fname, label) in enumerate(items[:cols * rows]):
        r, c = divmod(i, cols)
        x = MARGIN + c * (cell_w + gap)
        yy = y + r * (cell_h + gap)
        place_image(slide, ROOT / FIG_DIR_NAME / fname, x, yy, cell_w, cell_h - lab_h)
        tf = txbox(slide, x, yy + cell_h - lab_h, cell_w, lab_h)
        write(tf, label, size=8.5, colour=MUTED, align=PP_ALIGN.CENTER)
    footer(slide, d.get("source"), page)
    notes(slide, d["notes"])


# --- slide kinds for the methodological deck --------------------------------------------------
# These seven were added for the version a co-author reviews rather than the one that gets
# projected. That deck has to carry the chain, the parameters and the provenance, and doing that
# with the original kinds meant either a wall of text on a `figure` caption or an image of a table.
# Native shapes stay editable, reflow when a value changes, and print legibly.

ACCENTS = [BLUE, ORANGE, RGBColor(0x4D, 0x9A, 0x53), RGBColor(0x7B, 0x5E, 0xA7),
           RGBColor(0xC1, 0x62, 0x2E), RGBColor(0x3D, 0x6B, 0x7D)]
PANEL = RGBColor(0xF7, 0xF9, 0xFB)
ZEBRA = RGBColor(0xF2, 0xF5, 0xF8)


def est_lines(text, width, size_pt, factor=0.0072):
    """How many lines `text` will wrap to inside a box `width` EMU wide, set at `size_pt`.

    Segoe UI runs close to half an em per character over mixed-case English prose, so a character
    is roughly size_pt/144 inches. The factor is deliberately a shade generous: under-estimating
    clips text off the slide, over-estimating only leaves white space.
    """
    if not text:
        return 1
    per_line = max(8, int(width / Inches(factor * size_pt)))
    return sum(max(1, -(-len(ln) // per_line)) for ln in str(text).split("\n"))


def est_h(text, width, size_pt, line=1.24, pad=0.10):
    """Height an estimated block of text needs, with a little padding."""
    return Inches(est_lines(text, width, size_pt) * size_pt * line / 72 + pad)


def _body_lines(tf, lines, *, size, colour=MUTED, space_after=7, line=1.16):
    """Write a list of strings as consecutive paragraphs in one frame."""
    for j, ln in enumerate(lines):
        if j == 0:
            write(tf, ln, size=size, colour=colour, para=0, space_after=space_after, line=line)
        else:
            para(tf, ln, size=size, colour=colour, space_after=space_after, line=line)


def s_stepper(prs, d, page):
    """The chain as numbered steps across the page, each carrying its own parameters.

    The point of this kind is that a reader can see the order of operations and the value applied at
    each one without holding the previous slide in their head, which is exactly the complaint a
    method section usually earns.
    """
    slide = blank(prs)
    y = title_block(slide, d["title"])
    steps = d["steps"]
    n = len(steps)
    gap = Inches(0.30)
    col_w = int((CONTENT_W - gap * (n - 1)) / n)
    foot = d.get("foot")
    foot_h = (est_h(foot, CONTENT_W, BODY_PT - 1) + Inches(0.22)) if foot else Inches(0.1)
    card_h = H - y - Inches(0.72) - foot_h
    # A seven-step chain leaves each column about 1.5 in wide, and "Interpolate" set at 15.5 pt
    # breaks across two lines mid-word. The heading follows the column rather than the other way
    # round, because shortening the words would cost more than the point size does.
    inner = col_w - Inches(0.52)
    head_pt = 15.5 if inner >= Inches(1.65) else 14 if inner >= Inches(1.35) else \
        12.5 if inner >= Inches(1.05) else 11
    body_pt = 11.5 if inner >= Inches(1.35) else 10.5
    for i, st in enumerate(steps):
        x = MARGIN + i * (col_w + gap)
        accent = ACCENTS[i % len(ACCENTS)]
        rect(slide, x, y, col_w, card_h, PANEL)
        rect(slide, x, y, col_w, Inches(0.075), accent)
        # The step number sits in the corner rather than in the flow, so the eye reads the heading
        # first and uses the number only to keep its place.
        tf = txbox(slide, x + Inches(0.26), y + Inches(0.26), Inches(0.5), Inches(0.34))
        write(tf, str(i + 1), size=15, bold=True, colour=accent)
        tf = txbox(slide, x + Inches(0.26), y + Inches(0.66), inner, Inches(0.62))
        write(tf, st["head"], size=head_pt, bold=True, line=1.06)
        tf = txbox(slide, x + Inches(0.26), y + Inches(1.36), inner, card_h - Inches(2.15))
        _body_lines(tf, st["body"], size=body_pt)
        if st.get("param"):
            chip = rect(slide, x + Inches(0.26), y + card_h - Inches(0.62),
                        col_w - Inches(0.52), Inches(0.4), WHITE, line=RULE)
            ctf = chip.text_frame
            ctf.margin_left = ctf.margin_right = Inches(0.09)
            ctf.margin_top = ctf.margin_bottom = 0
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            write(ctf, st["param"], size=10, colour=INK, align=PP_ALIGN.CENTER)
        if i < n - 1:
            tf = txbox(slide, x + col_w, y + card_h / 2 - Inches(0.22), gap, Inches(0.44))
            write(tf, "›", size=22, colour=FAINT, align=PP_ALIGN.CENTER)
    if foot:
        tf = txbox(slide, MARGIN, y + card_h + Inches(0.22), CONTENT_W, foot_h)
        write(tf, foot, size=BODY_PT - 1, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_params(prs, d, page):
    """A parameter table with the file and line each value is set at.

    The provenance column is the reason this kind exists: a co-author asked to trust a number should
    be able to go and look at it, and a slide that says "IDW, 50 km" without saying where 50 km is
    written is asking to be taken on faith.
    """
    slide = blank(prs)
    y = title_block(slide, d["title"])
    rows = d["rows"]
    foot = d.get("foot")
    k_w = int(CONTENT_W * d.get("key_frac", 0.24))
    s_w = int(CONTENT_W * d.get("src_frac", 0.24))
    v_w = CONTENT_W - k_w - s_w
    foot_h = (est_h(foot, CONTENT_W, BODY_PT - 2) + Inches(0.18)) if foot else Inches(0)
    avail = H - y - Inches(0.72) - foot_h

    # Rows are measured, not assumed. A fixed row height clipped every value that wrapped, and on
    # the open-questions slide that was three rows out of nine. If the measured table overruns the
    # slide the type comes down rather than the text being cut.
    fs = d.get("size", 12)
    while True:
        heights = [max(Inches(0.30),
                       est_h(r[1], v_w - Inches(0.2), fs, line=1.22, pad=0.12)) for r in rows]
        if sum(heights) <= avail or fs <= 9:
            break
        fs -= 0.5

    ry = y
    for i, (r, row_h) in enumerate(zip(rows, heights)):
        if i % 2 == 0:
            rect(slide, MARGIN, ry, CONTENT_W, row_h, ZEBRA)
        pad = Inches(0.06)
        tf = txbox(slide, MARGIN + Inches(0.16), ry + pad, k_w - Inches(0.2), row_h)
        write(tf, r[0], size=fs, bold=True, colour=RGBColor(0x2C, 0x4A, 0x5E), line=1.22)
        tf = txbox(slide, MARGIN + k_w, ry + pad, v_w - Inches(0.2), row_h)
        write(tf, r[1], size=fs, line=1.22)
        if len(r) > 2 and r[2]:
            tf = txbox(slide, MARGIN + k_w + v_w, ry + pad, s_w, row_h)
            write(tf, r[2], size=fs - 3, colour=FAINT, align=PP_ALIGN.RIGHT)
        ry += row_h
    if foot:
        tf = txbox(slide, MARGIN, ry + Inches(0.18), CONTENT_W, foot_h)
        write(tf, foot, size=BODY_PT - 2, colour=MUTED, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_twocol(prs, d, page):
    """A figure beside headed blocks of explanation, rather than beside loose sentences.

    `figure_side` gives the reader a stack of equal-weight lines. When the explanation has structure
    (what it shows, how to read it, what it does not settle) the headings carry that structure and
    the reader can skip to the part they want.
    """
    slide = blank(prs)
    y = title_block(slide, d["title"])
    box_h = H - y - Inches(0.72)
    frac = d.get("image_frac", 0.52)
    img_w = int(CONTENT_W * frac)
    on_right = d.get("image_side", "left") == "right"
    img_x = (W - MARGIN - img_w) if on_right else MARGIN
    tx = MARGIN if on_right else (MARGIN + img_w + Inches(0.45))
    tw = CONTENT_W - img_w - Inches(0.45)
    place_image(slide, ROOT / d["image"], img_x, y, img_w, box_h)

    # The column is measured before anything is drawn. Blocks vary from one line to five, so a
    # fixed slot either clips the long ones or leaves holes under the short ones; and when the
    # whole column overruns the slide the type comes down rather than the last block falling off
    # the bottom, which is what happened as soon as a correction made one block longer.
    blocks = [dict(head=b["head"],
                   body=b["body"] if isinstance(b["body"], list) else [b["body"]])
              for b in d["blocks"]]
    head_pt, body_pt = 13.5, 12.0
    gap, lead = Inches(0.20), Inches(0.30)
    while True:
        heights = [sum(est_h(ln, tw, body_pt, line=1.18, pad=0.04) for ln in b["body"])
                   for b in blocks]
        total = sum(heights) + len(blocks) * (lead + gap)
        if total <= box_h or body_pt <= 9.5:
            break
        head_pt -= 0.5
        body_pt -= 0.5

    yy = y + Inches(0.06)
    for blk, h in zip(blocks, heights):
        tf = txbox(slide, tx, yy, tw, lead)
        write(tf, blk["head"], size=head_pt, bold=True, colour=RGBColor(0x2C, 0x4A, 0x5E))
        yy += lead
        tf = txbox(slide, tx, yy, tw, h)
        _body_lines(tf, blk["body"], size=body_pt, space_after=5, line=1.18)
        yy += h + gap
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


# --- code references ---------------------------------------------------------------------------
#
# The two parameter tables say where in the source each parameter is set. That used to be a string
# typed into the content file, "15_...:318", and a column of those reads as digits and ellipses
# rather than as a reference to anything.
#
# It is structured now: a row gives a file and the lines inside it, and the renderer writes the
# cell. Two things follow. The line numbers leave the slide, because a reviewer wants to know which
# file a parameter lives in and almost nobody wants line 271; they are kept in the data, checked at
# build time, and written into the speaker notes. And the file name is printed only where it
# changes, so a run of six rows from one script reads as one named block instead of six repetitions
# of a 28-character name.

SCRIPT_DIR = ROOT / "01_scripts"
REF_COUNT = [0]


def _loc_text(specs):
    """The line part of a reference: "318", "116, 339", "4-5"."""
    return ", ".join(f"{s[0]}-{s[1]}" if isinstance(s, tuple) else str(s) for s in specs)


def _loc_lines(specs):
    out = []
    for s in specs:
        out.extend(range(s[0], s[1] + 1) if isinstance(s, tuple) else [s])
    return out


def verify_ref(name, specs):
    """Stop the build if a reference no longer resolves to a real line of a real file.

    This catches the two ways a reference rots on its own: the script is renamed or moved, and the
    script is shortened past the line. It cannot catch a line that still exists but now holds
    something else, which is why the notes carry the numbers where a reader can check them.
    """
    path = SCRIPT_DIR / name
    if not path.exists():
        sys.exit(f"code reference {name}:{_loc_text(specs)} names a file that is not in "
                 f"{SCRIPT_DIR}")
    with open(path, encoding="utf-8", errors="replace") as fh:
        lines = fh.read().split("\n")
    wanted = _loc_lines(specs)
    past = [n for n in wanted if n < 1 or n > len(lines)]
    if past:
        sys.exit(f"code reference {name}:{_loc_text(specs)} points past the end of the file, "
                 f"which has {len(lines)} lines: {past}")
    blank = [n for n in wanted if not lines[n - 1].strip()]
    if blank:
        sys.exit(f"code reference {name}:{_loc_text(specs)} points at a blank line: {blank}\n"
                 f"  the file has almost certainly been edited since the reference was written")
    REF_COUNT[0] += len(wanted)


def resolve_refs(rows):
    """Turn the reference cells into display text, printing a file name only when it changes."""
    out, last = [], None
    for r in rows:
        ref = r[-1]
        if not isinstance(ref, tuple):
            out.append(list(r))
            continue
        name, specs = ref
        verify_ref(name, specs)
        out.append(list(r[:-1]) + [name if name != last else ""])
        last = name
    return out


def ref_notes(rows):
    """The exact lines, for the speaker notes, since they no longer appear on the slide."""
    parts = [f"{r[0]} = {r[-1][0]}:{_loc_text(r[-1][1])}"
             for r in rows if isinstance(r[-1], tuple)]
    return "Exact lines, kept off the slide: " + "; ".join(parts) + "." if parts else ""


def s_table(prs, d, page):
    """A real table in native shapes, with one row or column allowed to carry emphasis."""
    slide = blank(prs)
    y = head_of(slide, d)
    head, rows = d["head"], d["rows"]
    extra_note = ""
    if any(isinstance(r[-1], tuple) for r in rows):
        extra_note = ref_notes(rows)
        rows = resolve_refs(rows)
        d = dict(d, foot=((d.get("foot", "") + " ") if d.get("foot") else "")
                 + "A blank means the same file as the row above.")
    widths = d.get("widths") or [1.0 / len(head)] * len(head)
    cols = [int(CONTENT_W * f) for f in widths]
    xs, acc = [], MARGIN
    for c in cols:
        xs.append(acc)
        acc += c
    foot = d.get("foot")
    foot_h = (est_h(foot, CONTENT_W, BODY_PT - 2) + Inches(0.30)) if foot else Inches(0)
    avail = H - y - Inches(0.72) - foot_h
    emph = d.get("emphasis")

    # A row is as tall as its tallest cell needs. A fixed height suits a table of short values and
    # clips the moment one column carries a sentence, which is what a "why" column always does.
    fs = d.get("size", 12)
    head_h = Inches(0.40)
    floor = min(Inches(0.34), int((avail - head_h) / max(len(rows), 1)))
    while True:
        heights = [max(floor,
                       Inches(max(est_lines(str(c), cols[j] - Inches(0.24), fs)
                                  for j, c in enumerate(r)) * fs * 1.24 / 72 + 0.14))
                   for r in rows]
        if sum(heights) + head_h <= avail or fs <= 9:
            break
        fs -= 0.5

    if d.get("slim"):
        y = centred(y, avail, head_h + sum(heights))
    rect(slide, MARGIN, y, CONTENT_W, head_h, RGBColor(0x2C, 0x4A, 0x5E))
    for j, htxt in enumerate(head):
        al = PP_ALIGN.RIGHT if j and d.get("numeric", True) else PP_ALIGN.LEFT
        tf = txbox(slide, xs[j] + Inches(0.12), y + Inches(0.09), cols[j] - Inches(0.24), head_h)
        write(tf, htxt, size=fs - 1, bold=True, colour=WHITE, align=al)

    ry = y + head_h
    for i, (r, row_h) in enumerate(zip(rows, heights)):
        if emph is not None and i == emph:
            rect(slide, MARGIN, ry, CONTENT_W, row_h, RGBColor(0xE8, 0xF0, 0xF6))
        elif i % 2 == 0:
            rect(slide, MARGIN, ry, CONTENT_W, row_h, ZEBRA)
        for j, cell in enumerate(r):
            al = PP_ALIGN.RIGHT if j and d.get("numeric", True) else PP_ALIGN.LEFT
            bold = (emph is not None and i == emph) or (j == 0 and d.get("bold_first", True))
            tf = txbox(slide, xs[j] + Inches(0.12), ry + Inches(0.07), cols[j] - Inches(0.24),
                       row_h)
            write(tf, str(cell), size=fs, bold=bold, align=al, line=1.22)
        ry += row_h
    if foot:
        tf = txbox(slide, MARGIN, ry + Inches(0.18), CONTENT_W, foot_h)
        write(tf, foot, size=BODY_PT - 2, colour=MUTED, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, "\n\n".join(x for x in (d.get("notes", ""), extra_note) if x))


def s_datacard(prs, d, page):
    """One card per data source, each with its real dimensions as key-value rows."""
    slide = blank(prs)
    y = title_block(slide, d["title"])
    items = d["items"]
    n = len(items)
    gap = Inches(0.30)
    col_w = int((CONTENT_W - gap * (n - 1)) / n)
    foot = d.get("foot")
    card_h = H - y - Inches(0.72) - (Inches(0.82) if foot else Inches(0.1))
    for i, it in enumerate(items):
        x = MARGIN + i * (col_w + gap)
        accent = ACCENTS[i % len(ACCENTS)]
        rect(slide, x, y, col_w, card_h, PANEL)
        rect(slide, x, y, col_w, Inches(0.075), accent)
        tf = txbox(slide, x + Inches(0.26), y + Inches(0.28), col_w - Inches(0.52), Inches(0.6))
        write(tf, it["head"], size=15, bold=True, line=1.06)
        inner = col_w - Inches(0.52)
        sub_h = est_h(it["sub"], inner, 10.5, line=1.12, pad=0.06)
        tf = txbox(slide, x + Inches(0.26), y + Inches(0.92), inner, sub_h)
        write(tf, it["sub"], size=10.5, colour=MUTED, line=1.1)
        note_h = (est_h(it["note"], inner - Inches(0.2), 9.5, line=1.14, pad=0.16)
                  if it.get("note") else Inches(0))
        ry = y + Inches(0.98) + sub_h
        room = (y + card_h - note_h - Inches(0.22)) - ry
        pitch = min(Inches(0.36), int(room / max(len(it["rows"]), 1)))
        for k, v in it["rows"]:
            rect(slide, x + Inches(0.26), ry, col_w - Inches(0.52), Emu(9525), RULE)
            tf = txbox(slide, x + Inches(0.26), ry + Inches(0.08), int((col_w - Inches(0.52)) * 0.38),
                       Inches(0.28))
            write(tf, k, size=10, colour=MUTED)
            tf = txbox(slide, x + Inches(0.26) + int((col_w - Inches(0.52)) * 0.38), ry + Inches(0.08),
                       int((col_w - Inches(0.52)) * 0.62), Inches(0.28))
            write(tf, v, size=10, bold=True, align=PP_ALIGN.RIGHT)
            ry += pitch
        if it.get("note"):
            nb = rect(slide, x + Inches(0.26), y + card_h - note_h - Inches(0.14), inner,
                      note_h, RGBColor(0xFD, 0xF3, 0xE3))
            ntf = nb.text_frame
            ntf.margin_left = ntf.margin_right = Inches(0.1)
            ntf.margin_top = ntf.margin_bottom = Inches(0.05)
            ntf.word_wrap = True
            write(ntf, it["note"], size=9.5, colour=RGBColor(0x6B, 0x4C, 0x17), line=1.12)
    if foot:
        tf = txbox(slide, MARGIN, y + card_h + Inches(0.2), CONTENT_W, Inches(0.66))
        write(tf, foot, size=BODY_PT - 1, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_bignum(prs, d, page):
    """A row of headline figures. Used where the argument is the size of a number, not its shape."""
    slide = blank(prs)
    y = head_of(slide, d)
    items = d["items"]
    n = len(items)
    gap = Inches(0.28)
    col_w = int((CONTENT_W - gap * (n - 1)) / n)
    inner = col_w - Inches(0.5)
    lead = d.get("lead")
    if lead:
        lead_h = est_h(lead, CONTENT_W, BODY_PT, line=1.2)
        tf = txbox(slide, MARGIN, y, CONTENT_W, lead_h)
        write(tf, lead, size=BODY_PT, colour=MUTED, line=1.2)
        y += lead_h + Inches(0.18)

    # The headline is the size of the number, so the number must not wrap. "229,676 km2" set at
    # 34 pt needs 2.4 in and a four-column row gives it 2.3, which broke the unit onto a second
    # line and dropped it on top of the label. The type comes down until every value fits one line.
    val_pt = 34
    while val_pt > 17 and any(est_lines(it["value"], inner, val_pt) > 1 for it in items):
        val_pt -= 1.5
    val_h = Inches(val_pt * 1.06 / 72 + 0.08)
    lab_h = max(est_h(it["label"], inner, 11.5, line=1.16, pad=0.04) for it in items)
    box_h = Inches(0.24) + val_h + Inches(0.10) + lab_h + Inches(0.18)
    if d.get("slim"):
        body_h = est_h("\n".join(d.get("body") or []), CONTENT_W, BODY_PT - 1, line=1.22)
        y = centred(y, H - y - Inches(0.72), box_h + Inches(0.34) + body_h)

    for i, it in enumerate(items):
        x = MARGIN + i * (col_w + gap)
        accent = it.get("accent") or ACCENTS[i % len(ACCENTS)]
        rect(slide, x, y, col_w, box_h, PANEL)
        rect(slide, x, y, Inches(0.075), box_h, accent)
        tf = txbox(slide, x + Inches(0.32), y + Inches(0.24), inner, val_h)
        write(tf, it["value"], size=val_pt, bold=True, colour=accent, line=0.98)
        tf = txbox(slide, x + Inches(0.32), y + Inches(0.34) + val_h, inner, lab_h)
        write(tf, it["label"], size=11.5, colour=MUTED, line=1.14)
    body = d.get("body") or []
    if body:
        tf = txbox(slide, MARGIN, y + box_h + Inches(0.34), CONTENT_W, H - y - box_h - Inches(1.1))
        _body_lines(tf, body, size=BODY_PT - 1, colour=INK, space_after=11, line=1.22)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_annotated(prs, d, page):
    """A figure with numbered callouts pinned to it.

    Positions are given as fractions of the placed image, not as inches, so a figure that is
    regenerated at a different aspect ratio keeps its annotations on the right part of the picture.
    """
    slide = blank(prs)
    y = title_block(slide, d["title"])
    notes_w = int(CONTENT_W * d.get("notes_frac", 0.27))
    img_w = CONTENT_W - notes_w - Inches(0.4)
    box_h = H - y - Inches(0.72)
    left, top, w, h = place_image(slide, ROOT / d["image"], MARGIN, y, img_w, box_h)
    tx = MARGIN + img_w + Inches(0.4)
    yy = y + Inches(0.06)
    for i, c in enumerate(d["callouts"], start=1):
        accent = ACCENTS[(i - 1) % len(ACCENTS)]
        cx = left + int(w * c["at"][0])
        cy = top + int(h * c["at"][1])
        badge = rect(slide, cx - Inches(0.13), cy - Inches(0.13), Inches(0.26), Inches(0.26), accent)
        btf = badge.text_frame
        btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = 0
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(btf, str(i), size=10, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        tf = txbox(slide, tx, yy, Inches(0.26), Inches(0.3))
        write(tf, str(i), size=12, bold=True, colour=accent)
        tf = txbox(slide, tx + Inches(0.3), yy, notes_w - Inches(0.3), Inches(1.0))
        write(tf, c["text"], size=11.5, colour=INK, line=1.16)
        est = 1 + len(c["text"]) // int((notes_w - Inches(0.3)) / Inches(0.075))
        yy += Inches(0.20) * est + Inches(0.22)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


# --- v4 layout: the picture gets the page, the prose becomes a drawing ------------------------
#
# Measured on the v3 deck: 33,860 characters of prose over 54 slides, and the largest figure took
# 56% of the page while the two national maps took 20%. The wording was only half the problem. A
# two-line title block plus a caption band spend 3.3 in of a 7.5 in page before the picture is
# placed, so a near-square map can never exceed 4.7 in wide however short the text gets. The kinds
# below change the geometry: a one-line title, no caption band, the border cropped off the PNG,
# and a side rail carrying the legend so the map itself keeps the whole height. That side rail is
# what the supervision meeting of 24 August 2026 asked for in as many words.
#
# They are new kinds rather than edits to the existing ones because slides() and annex() still
# build the conference talk and its backup from the same file, and those two decks are finished.

TRIM_DIR = ROOT / "02_outputs" / "_deck_trimmed"


def trimmed(path, crop=None):
    """A copy of `path` with its uniform border removed, cached beside the outputs.

    R draws onto a canvas of fixed aspect, so a portrait map inside a landscape device arrives with
    the unused canvas as white pixels: fig20_15 is 27% blank by width, which is 27% of the slide
    the picture cannot use. `crop` additionally removes a fraction of the top and bottom, which is
    how the map's own printed title and legend row come off when the slide already carries both.

    Cropping happens here and not in the R scripts because the same PNG is a figure of the method
    book, where the printed title and the margin are wanted. The fractions are therefore tied to
    the layout of 19_cropland_viability_national.R and are passed per slide, next to the image.
    """
    if path.suffix.lower() not in (".png", ".jpg", ".jpeg"):
        return path                                  # a GIF must reach PowerPoint untouched
    tag = "" if not crop else "_c%02d%02d" % (round(crop[0] * 100), round(crop[1] * 100))
    out = TRIM_DIR / f"{path.stem}{tag}{path.suffix}"
    if out.exists() and out.stat().st_mtime >= path.stat().st_mtime:
        return out
    TRIM_DIR.mkdir(parents=True, exist_ok=True)
    with Image.open(path) as im:
        im = im.convert("RGB")
        w, h = im.size
        if crop:
            im = im.crop((0, int(h * crop[0]), w, h - int(h * crop[1])))
        bg = Image.new("RGB", im.size, im.getpixel((0, 0)))
        box = ImageChops.difference(im, bg).getbbox()
        if box:
            pad = max(2, int(0.004 * max(im.size)))
            im = im.crop((max(0, box[0] - pad), max(0, box[1] - pad),
                          min(im.size[0], box[2] + pad), min(im.size[1], box[3] + pad)))
        im.save(out)
    return out


def figure_of(d):
    """The file a slide's picture should actually be built from, trimmed and optionally cropped."""
    return trimmed(ROOT / d["image"], d.get("crop"))


def head(slide, d, *, slim=True):
    """The slide's head: a title that names the subject, and optionally the finding under it.

    Keeping the two apart is the whole point. `title` says what the slide is about, so a reader
    scanning the deck can navigate it, and `kicker` says what to take away, so a reader who stops
    on this slide does not have to work it out from the evidence. Written as one sentence in the
    title, as the earlier decks did, the two jobs fight each other and the title stops being a
    title.
    """
    y = title_slim(slide, d["title"]) if slim else title_block(slide, d["title"])
    kick = d.get("kicker")
    if kick:
        h = est_h(kick, CONTENT_W, 13.5, line=1.18)
        tf = txbox(slide, MARGIN, y - Inches(0.02), CONTENT_W, h)
        write(tf, kick, size=13.5, colour=MUTED, line=1.18)
        y += h + Inches(0.12)
    return y


def head_of(slide, d):
    """The title treatment a slide asks for: slim for v4, the taller block for the rest."""
    return head(slide, d, slim=bool(d.get("slim")))


def title_slim(slide, text, *, size=21, y=None):
    """A one-line title and its hairline, in half the height of title_block().

    title_block() reserves two lines whenever the title runs past 62 characters, and every
    assertion title does. Here the type shrinks instead, which costs nothing at the back of a room
    and returns an inch and a half of page to the figure.
    """
    y = Inches(0.34) if y is None else y
    while size > 15 and est_lines(text, CONTENT_W, size) > 1:
        size -= 1
    h = Inches(size * 1.15 / 72 + 0.06)
    tf = txbox(slide, MARGIN, y, CONTENT_W, h)
    write(tf, text, size=size, bold=True, line=0.95)
    ry = y + h + Inches(0.10)
    rect(slide, MARGIN, ry, CONTENT_W, Emu(9525), RULE)
    return ry + Inches(0.16)


def rrect(slide, x, y, w, h, fill, *, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    """rect() with a rounded corner and its adjustment pulled back to something restrained."""
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    try:
        sh.adjustments[0] = 0.07
    except (IndexError, ValueError):
        pass
    return sh


def centred(y, avail, block_h):
    """Where to start a block of `block_h` so it sits centred in the band under the title.

    Sizing a panel to the space rather than to its contents is what left the first build's step
    cards two thirds empty, with the parameter chip stranded at the bottom of a tall white box. The
    panels are measured from their text now, which means they no longer reach the footer, so the
    remainder is split above and below instead of all falling underneath.
    """
    return y + max(Emu(0), int((avail - block_h) * 0.38))


def chip(slide, x, y, w, text, *, size=9, fill=PANEL, colour=MUTED, bold=False, h=None):
    """A small tinted label. Used for parameters, units and provenance under a heading."""
    h = Inches(size * 1.5 / 72 + 0.09) if h is None else h
    sh = rrect(slide, x, y, w, h, fill)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(tf, text, size=size, colour=colour, bold=bold, align=PP_ALIGN.CENTER)
    return y + h


def s_figure_max(prs, d, page):
    """One picture and its assertion. The picture takes whatever the title and caption leave."""
    slide = blank(prs)
    y = head(slide, d)
    cap = d.get("caption")
    cap_h = est_h(cap, CONTENT_W, CAP_PT - 1, line=1.16, pad=0.10) if cap else Inches(0)
    box_h = H - y - Inches(0.64) - cap_h
    left, top, w, h = place_image(slide, figure_of(d), MARGIN, y, CONTENT_W, box_h)
    if cap:
        tf = txbox(slide, MARGIN, y + box_h + Inches(0.04), CONTENT_W, cap_h)
        write(tf, cap, size=CAP_PT - 1, colour=MUTED, line=1.16)
    if d.get("gif"):
        badge = rrect(slide, left + w - Inches(1.6), top + Inches(0.06), Inches(1.5),
                      Inches(0.26), INK)
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(tf, "animated · press F5", size=9, colour=WHITE, align=PP_ALIGN.CENTER)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_map(prs, d, page):
    """A map at full page height with its legend and its headline figures in a rail beside it.

    The three colours are the result, so the rail states what each one means and how much of the
    country it covers in the same row. That merges the legend with the numbers that used to sit in
    a caption under a picture too small to read.
    """
    slide = blank(prs)
    rail_w = Inches(3.5)
    # The rule follows the title rather than sitting at a fixed height. Titles became short names
    # of the slide's subject, and a one-line title left an inch of nothing above the legend.
    t_w = rail_w - Inches(0.2)
    t_h = est_h(d["title"], t_w, 20, line=0.98, pad=0.06)
    tf = txbox(slide, MARGIN, Inches(0.42), t_w, t_h)
    write(tf, d["title"], size=20, bold=True, line=0.98)
    rect(slide, MARGIN, Inches(0.42) + t_h + Inches(0.30), t_w, Emu(9525), RULE)

    # The legend keys are the three classes, so the content file never names a colour: blue for
    # both cultivars, orange for the mutant alone and red for neither are fixed by 00_map_layout.R
    # and the same three appear on every figure in the project.
    by_key = {"both": BLUE, "only": ORANGE, "neither": RED}
    yy = Inches(0.42) + t_h + Inches(0.56)
    for row in d["legend"]:
        colour = row.get("colour") or by_key[row["key"]]
        row = dict(row, colour=colour)
        rect(slide, MARGIN, yy + Inches(0.05), Inches(0.24), Inches(0.24), row["colour"])
        tf = txbox(slide, MARGIN + Inches(0.36), yy, rail_w - Inches(0.6), Inches(0.34))
        write(tf, row["label"], size=12.5, colour=INK, line=1.0)
        if row.get("value"):
            tf = txbox(slide, MARGIN + Inches(0.36), yy + Inches(0.28),
                       rail_w - Inches(0.6), Inches(0.4))
            write(tf, row["value"], size=19, bold=True, colour=row["colour"], line=0.98)
            yy += Inches(0.80)
        else:
            yy += Inches(0.42)

    for extra in d.get("rail", []):
        tf = txbox(slide, MARGIN, yy + Inches(0.16), rail_w - Inches(0.2),
                   est_h(extra, rail_w - Inches(0.2), 11, line=1.18))
        write(tf, extra, size=11, colour=MUTED, line=1.18)
        yy += Inches(0.16) + est_h(extra, rail_w - Inches(0.2), 11, line=1.18)

    x = MARGIN + rail_w + Inches(0.24)
    place_image(slide, figure_of(d), x, Inches(0.30), W - x - MARGIN, H - Inches(0.74))
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_figure_note(prs, d, page):
    """A picture with at most three short notes, beside it or underneath it.

    The v3 twocol gave the picture 50% of the width and four paragraphs the rest. Here the picture
    takes what it can use and the notes are captions rather than argument: whatever needs a
    paragraph belongs in the method book, which travels with this deck.

    Which side the notes go on is decided by the picture, not by hand. A note column 30% wide
    leaves a box of roughly 8.5 by 5.9 in, so any figure wider than about 1.6:1 runs out of width
    before it runs out of height and the slide ends up with a strip of white under it. Those go
    underneath instead, where the figure gets the full page width. Deciding it here rather than in
    the content file means a figure regenerated at a new aspect lands in the right layout on its
    own, which four of them did not in the first build.
    """
    slide = blank(prs)
    y = head(slide, d)
    path = figure_of(d)
    with Image.open(path) as im:
        aspect = im.size[0] / im.size[1]
    side = aspect <= d.get("side_max_aspect", 1.6)
    ns = d["notes_side"]

    if side:
        note_w = int(CONTENT_W * d.get("note_frac", 0.29))
        img_w = CONTENT_W - note_w - Inches(0.34)
        place_image(slide, path, MARGIN, y, img_w, H - y - Inches(0.64))
        x = MARGIN + img_w + Inches(0.34)
        yy = y + Inches(0.06)
        for i, n in enumerate(ns):
            accent = ACCENTS[i % len(ACCENTS)]
            rect(slide, x, yy + Inches(0.04), Inches(0.06), Inches(0.22), accent)
            tf = txbox(slide, x + Inches(0.18), yy, note_w - Inches(0.18), Inches(0.3))
            write(tf, n["head"], size=12.5, bold=True, line=1.0)
            body_h = est_h(n["body"], note_w - Inches(0.18), 11.5, line=1.2)
            tf = txbox(slide, x + Inches(0.18), yy + Inches(0.30), note_w - Inches(0.18), body_h)
            write(tf, n["body"], size=11.5, colour=MUTED, line=1.2)
            yy += Inches(0.30) + body_h + Inches(0.26)
    else:
        gap = Inches(0.30)
        col_w = int((CONTENT_W - gap * (len(ns) - 1)) / len(ns))
        row_h = Inches(0.30) + max(est_h(n["body"], col_w, 11.5, line=1.2) for n in ns)
        # The footer sits 0.5 in from the bottom, so the note row has to clear it: at 0.52
        # the second line of a two-line note lands on the source line.
        place_image(slide, path, MARGIN, y, CONTENT_W, H - y - row_h - Inches(0.88))
        ny = H - Inches(0.76) - row_h
        for i, n in enumerate(ns):
            x = MARGIN + i * (col_w + gap)
            accent = ACCENTS[i % len(ACCENTS)]
            rect(slide, x, ny + Inches(0.04), Inches(0.06), Inches(0.22), accent)
            tf = txbox(slide, x + Inches(0.18), ny, col_w - Inches(0.18), Inches(0.3))
            write(tf, n["head"], size=12.5, bold=True, line=1.0)
            tf = txbox(slide, x + Inches(0.18), ny + Inches(0.30), col_w - Inches(0.18),
                       row_h - Inches(0.30))
            write(tf, n["body"], size=11.5, colour=MUTED, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_flow(prs, d, page):
    """Numbered steps in a row with arrows between them, and a parameter chip under each.

    Replaces the seven-step list, which read as prose set in columns. A chain is a shape, and
    drawing it as one lets the parameter that governs each step sit on the step itself.
    """
    slide = blank(prs)
    y = head(slide, d)
    steps = d["steps"]
    n = len(steps)
    arrow = Inches(0.20)
    col_w = int((CONTENT_W - arrow * (n - 1)) / n)
    inner = col_w - Inches(0.28)
    foot = d.get("foot")
    foot_h = est_h(foot, CONTENT_W, 12, line=1.2) + Inches(0.24) if foot else Inches(0)
    avail = H - y - Inches(0.58) - foot_h
    head_pt = 14 if inner >= Inches(1.2) else 12
    body_pt = 12 if inner >= Inches(1.2) else 10
    head_h = max(est_h(st["head"], inner, head_pt, line=1.0) for st in steps)
    body_h = max(est_h(st["body"], inner, body_pt, line=1.2) for st in steps)
    chip_h = Inches(0.36) if any(st.get("param") for st in steps) else Inches(0)
    # A short chain leaves the band under the title half empty, so the card takes a floor of
    # 45% of it. Below that the seven blocks read as labels rather than as a pipeline.
    box_h = min(avail, max(Inches(0.46) + head_h + Inches(0.10) + body_h + Inches(0.24)
                           + chip_h + Inches(0.16), int(avail * 0.45)))
    y = centred(y, avail, box_h)

    for i, st in enumerate(steps):
        x = MARGIN + i * (col_w + arrow)
        accent = ACCENTS[i % len(ACCENTS)]
        rrect(slide, x, y, col_w, box_h, PANEL)
        rect(slide, x, y, col_w, Inches(0.07), accent)
        tf = txbox(slide, x + Inches(0.14), y + Inches(0.20), Inches(0.4), Inches(0.24))
        write(tf, str(i + 1), size=10, bold=True, colour=accent)
        tf = txbox(slide, x + Inches(0.14), y + Inches(0.46), inner, head_h)
        write(tf, st["head"], size=head_pt, bold=True, line=1.0)
        tf = txbox(slide, x + Inches(0.14), y + Inches(0.56) + head_h, inner, body_h)
        write(tf, st["body"], size=body_pt, colour=MUTED, line=1.2)
        if st.get("param"):
            chip(slide, x + Inches(0.14), y + box_h - Inches(0.52), inner, st["param"],
                 size=8.5, h=Inches(0.36), fill=WHITE)
        if i < n - 1:
            a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + col_w + Inches(0.03),
                                       y + box_h / 2 - Inches(0.07), Inches(0.14), Inches(0.14))
            a.fill.solid()
            a.fill.fore_color.rgb = RULE
            a.line.fill.background()
            a.shadow.inherit = False
    if foot:
        tf = txbox(slide, MARGIN, H - Inches(0.58) - foot_h, CONTENT_W, foot_h)
        write(tf, foot, size=12, colour=INK, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_lanes(prs, d, page):
    """Two or three labelled lanes with chips inside. Where a step runs, not what it does."""
    slide = blank(prs)
    y = head(slide, d)
    lanes = d["lanes"]
    foot = d.get("foot")
    foot_h = est_h(foot, CONTENT_W, 12, line=1.2) + Inches(0.24) if foot else Inches(0)
    avail = H - y - Inches(0.58) - foot_h
    gap = Inches(0.22)
    lab_w = Inches(2.05)
    # A lane is as tall as its tallest chip, not as tall as the space allows. Its own label needs
    # room too, so both are measured and the taller wins.
    item_w = min(int((CONTENT_W - lab_w - Inches(0.3) - Inches(0.16) * (len(ln["items"]) - 1))
                     / len(ln["items"])) for ln in lanes)
    need = Inches(0)
    for ln in lanes:
        body = max(est_h(it["body"], item_w - Inches(0.28), 10, line=1.16) for it in ln["items"])
        lab = est_h(ln["sub"], lab_w - Inches(0.3), 10.5, line=1.18) + Inches(0.66)
        need = max(need, Inches(0.40) + Inches(0.36) + body + Inches(0.48), lab + Inches(0.2))
    lane_h = min(int((avail - gap * (len(lanes) - 1)) / len(lanes)), need)
    y = centred(y, avail, lane_h * len(lanes) + gap * (len(lanes) - 1))
    for i, ln in enumerate(lanes):
        ly = y + i * (lane_h + gap)
        accent = ln.get("colour") or ACCENTS[i % len(ACCENTS)]
        rrect(slide, MARGIN, ly, CONTENT_W, lane_h, PANEL)
        rect(slide, MARGIN, ly, Inches(0.075), lane_h, accent)
        tf = txbox(slide, MARGIN + Inches(0.26), ly + Inches(0.22), lab_w - Inches(0.3),
                   Inches(0.42))
        write(tf, ln["head"], size=15, bold=True, colour=accent, line=1.0)
        tf = txbox(slide, MARGIN + Inches(0.26), ly + Inches(0.66), lab_w - Inches(0.3),
                   lane_h - Inches(0.8))
        write(tf, ln["sub"], size=10.5, colour=MUTED, line=1.18)
        items = ln["items"]
        iw = int((CONTENT_W - lab_w - Inches(0.3) - Inches(0.16) * (len(items) - 1)) / len(items))
        for j, it in enumerate(items):
            ix = MARGIN + lab_w + Inches(0.1) + j * (iw + Inches(0.16))
            rrect(slide, ix, ly + Inches(0.24), iw, lane_h - Inches(0.48), WHITE, line=RULE)
            tf = txbox(slide, ix + Inches(0.14), ly + Inches(0.40), iw - Inches(0.28),
                       Inches(0.36))
            write(tf, it["head"], size=11.5, bold=True, line=1.0)
            tf = txbox(slide, ix + Inches(0.14), ly + Inches(0.76), iw - Inches(0.28),
                       lane_h - Inches(1.06))
            write(tf, it["body"], size=10, colour=MUTED, line=1.16)
    if foot:
        tf = txbox(slide, MARGIN, H - Inches(0.58) - foot_h, CONTENT_W, foot_h)
        write(tf, foot, size=12, colour=INK, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_cards(prs, d, page):
    """Two to five cards, side by side or stacked as rows, each with an optional measured bar.

    The bar is drawn to scale from `bar`, so a card that says one dataset is fifteen times another
    shows it rather than asserting it. Rows are used where the cards are a numbered list, which is
    what the open questions are.
    """
    slide = blank(prs)
    y = head(slide, d)
    items = d["items"]
    foot = d.get("foot")
    foot_h = est_h(foot, CONTENT_W, 12, line=1.2) + Inches(0.22) if foot else Inches(0)
    avail = H - y - Inches(0.58) - foot_h

    if d.get("rows"):
        gap = Inches(0.14)
        body_w = CONTENT_W - Inches(3.0) - (Inches(2.5) if any(i.get("stat") for i in items)
                                            else Inches(0))
        row_h = min(int((avail - gap * (len(items) - 1)) / len(items)),
                    max(est_h(it["body"], body_w, 12, line=1.2) for it in items) + Inches(0.36))
        y = centred(y, avail, row_h * len(items) + gap * (len(items) - 1))
        for i, it in enumerate(items):
            ry = y + i * (row_h + gap)
            accent = it.get("colour") or ACCENTS[i % len(ACCENTS)]
            rrect(slide, MARGIN, ry, CONTENT_W, row_h, PANEL)
            rect(slide, MARGIN, ry, Inches(0.07), row_h, accent)
            tf = txbox(slide, MARGIN + Inches(0.26), ry + Inches(0.16), Inches(2.6),
                       row_h - Inches(0.3))
            write(tf, it["head"], size=13, bold=True, colour=accent, line=1.02)
            stat = it.get("stat")
            stat_w = Inches(2.5) if stat else Inches(0)
            bw = CONTENT_W - Inches(3.0) - stat_w
            tf = txbox(slide, MARGIN + Inches(2.95), ry + Inches(0.18), bw, row_h - Inches(0.32))
            write(tf, it["body"], size=12, colour=INK, line=1.2)
            if stat:
                tf = txbox(slide, W - MARGIN - stat_w, ry + Inches(0.22), stat_w - Inches(0.14),
                           Inches(0.4))
                write(tf, stat, size=11, colour=MUTED, align=PP_ALIGN.RIGHT, line=1.1)
        if foot:
            tf = txbox(slide, MARGIN, H - Inches(0.58) - foot_h, CONTENT_W, foot_h)
            write(tf, foot, size=12, colour=INK, line=1.2)
        footer(slide, d.get("source"), page)
        notes(slide, d.get("notes", ""))
        return

    n = len(items)
    gap = Inches(0.26)
    col_w = int((CONTENT_W - gap * (n - 1)) / n)
    inner = col_w - Inches(0.52)
    body_pt = 12.5 if n <= 3 else 11.5
    icon_h = Inches(1.24) if any(it.get("icon") for it in items) else Inches(0)
    card_h = min(avail,
                 Inches(0.30)
                 + max(est_h(it["head"], inner, 15, line=1.02) for it in items)
                 + Inches(0.10)
                 + (icon_h + Inches(0.14) if icon_h else Inches(0))
                 + (Inches(0.52) if any(it.get("stat") for it in items) else Inches(0))
                 + (Inches(0.30) if any(it.get("bar") is not None for it in items) else Inches(0))
                 + max(est_h(it["body"], inner, body_pt, line=1.2) for it in items)
                 + (Inches(0.56) if any(it.get("chip") for it in items) else Inches(0.16)))
    y = centred(y, avail, card_h)
    for i, it in enumerate(items):
        x = MARGIN + i * (col_w + gap)
        accent = it.get("colour") or ACCENTS[i % len(ACCENTS)]
        rrect(slide, x, y, col_w, card_h, PANEL)
        rect(slide, x, y, col_w, Inches(0.075), accent)
        tf = txbox(slide, x + Inches(0.26), y + Inches(0.30), inner, Inches(0.6))
        write(tf, it["head"], size=15, bold=True, line=1.02)
        yy = y + Inches(0.30) + est_h(it["head"], inner, 15, line=1.02) + Inches(0.10)
        if icon_h:
            # A card naming a dataset or a map is quicker to recognise by its picture than by its
            # heading, so the thumbnail sits directly under the heading. The band is reserved on
            # every card in the row, not only the ones carrying a picture: reserving it per card
            # would leave the headline figures at two different heights across the row, which reads
            # as a mistake rather than as a distinction. It is trimmed like any other figure, so a
            # white border cannot eat the band.
            if it.get("icon"):
                place_image(slide, trimmed(ROOT / it["icon"]), x + Inches(0.26), yy, inner, icon_h)
            yy += icon_h + Inches(0.14)
        if it.get("stat"):
            tf = txbox(slide, x + Inches(0.26), yy, inner, Inches(0.5))
            write(tf, it["stat"], size=25, bold=True, colour=accent, line=0.98)
            yy += Inches(0.52)
        if it.get("bar") is not None:
            rect(slide, x + Inches(0.26), yy, inner, Inches(0.12), RULE)
            rect(slide, x + Inches(0.26), yy, int(inner * max(0.02, it["bar"])), Inches(0.12),
                 accent)
            yy += Inches(0.30)
        tf = txbox(slide, x + Inches(0.26), yy, inner, card_h - (yy - y) - Inches(0.56))
        write(tf, it["body"], size=body_pt, colour=MUTED, line=1.2)
        if it.get("chip"):
            chip(slide, x + Inches(0.26), y + card_h - Inches(0.46), inner, it["chip"],
                 size=9.5, h=Inches(0.34), fill=WHITE)
    if foot:
        tf = txbox(slide, MARGIN, H - Inches(0.58) - foot_h, CONTENT_W, foot_h)
        write(tf, foot, size=12, colour=INK, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_timeline(prs, d, page):
    """Labelled bands on a shared year axis.

    The four analysis windows are a picture of a timeline, and were being described in a table with
    five columns. Drawing them to scale shows the property the table had to assert, which is that
    they tile the century without a gap or an overlap.
    """
    slide = blank(prs)
    y = head(slide, d)
    y0, y1 = d["span"]
    lab_w = Inches(2.3)
    ax_x = MARGIN + lab_w
    ax_w = CONTENT_W - lab_w - Inches(0.2)
    foot = d.get("foot")
    foot_h = est_h(foot, CONTENT_W, 12, line=1.2) + Inches(0.24) if foot else Inches(0)
    avail = H - y - Inches(0.58) - foot_h - Inches(0.78)

    def at(year):
        return ax_x + int(ax_w * (year - y0) / (y1 - y0))

    # A band is a span on an axis, so it stays a bar of fixed depth however much room the slide
    # has. Letting it grow to the row height turned the first build's four windows into four
    # rectangles the size of playing cards, which reads as a bar chart of nothing.
    bands = d["bands"]
    gap = Inches(0.12)
    band_h = Inches(0.50)
    row_h = min(int((avail - gap * (len(bands) - 1)) / len(bands)), Inches(0.98))
    block = row_h * len(bands) + gap * (len(bands) - 1)
    y = y + max(Emu(0), int((avail - block) / 2))
    for i, b in enumerate(bands):
        by = y + i * (row_h + gap)
        mid = by + row_h / 2
        accent = b.get("colour") or ACCENTS[i % len(ACCENTS)]
        tf = txbox(slide, MARGIN, mid - Inches(0.34), lab_w - Inches(0.18), Inches(0.3))
        write(tf, b["head"], size=12.5, bold=True, line=1.02)
        if b.get("sub"):
            tf = txbox(slide, MARGIN, mid - Inches(0.04), lab_w - Inches(0.18), Inches(0.3))
            write(tf, b["sub"], size=10, colour=MUTED, line=1.1)
        rect(slide, ax_x, mid - Emu(4763), ax_w, Emu(9525), RULE)
        x0, x1 = at(b["from"]), at(b["to"])
        sh = rrect(slide, x0, mid - band_h / 2, max(Inches(0.3), x1 - x0), band_h, accent)
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = Inches(0.08)
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(tf, b["label"], size=10.5, colour=WHITE, bold=True, align=PP_ALIGN.CENTER)

    ty = y + block + Inches(0.16)
    rect(slide, ax_x, ty, ax_w, Emu(9525), RULE)
    for year in d.get("ticks", []):
        rect(slide, at(year), ty, Emu(9525), Inches(0.09), FAINT)
        tf = txbox(slide, at(year) - Inches(0.3), ty + Inches(0.12), Inches(0.6), Inches(0.24))
        write(tf, str(year), size=9.5, colour=FAINT, align=PP_ALIGN.CENTER)
    if foot:
        tf = txbox(slide, MARGIN, H - Inches(0.58) - foot_h, CONTENT_W, foot_h)
        write(tf, foot, size=12, colour=INK, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_scale(prs, d, page):
    """A chill axis with the two cultivar requirements marked and values placed against them.

    Every classification in this work is a comparison against 47.5 and 33.7 chill portions, and
    saying so in a sentence makes the reader hold two numbers in their head. On an axis the three
    bands are visible at once, and a station or a scenario is a mark on it.
    """
    slide = blank(prs)
    y = head(slide, d)
    lo, hi = d["span"]
    ax_x, ax_w = MARGIN + Inches(0.4), CONTENT_W - Inches(0.8)
    bar_h = Inches(0.95)
    foot = d.get("foot")
    foot_h = est_h(foot, CONTENT_W, 13, line=1.24) if foot else Inches(0)
    avail = H - y - Inches(0.58)
    bar_y = centred(y, avail, Inches(1.30) + bar_h + Inches(0.70) + foot_h) + Inches(1.30)

    def at(v):
        return ax_x + int(ax_w * (v - lo) / (hi - lo))

    t_low, t_high = d["thresholds"]
    rect(slide, ax_x, bar_y, at(t_low) - ax_x, bar_h, RED)
    rect(slide, at(t_low), bar_y, at(t_high) - at(t_low), bar_h, ORANGE)
    rect(slide, at(t_high), bar_y, ax_x + ax_w - at(t_high), bar_h, BLUE)
    for lab, x0, x1 in (("neither", ax_x, at(t_low)),
                        ("only 'Búlida Precoz'", at(t_low), at(t_high)),
                        ("both cultivars", at(t_high), ax_x + ax_w)):
        tf = txbox(slide, x0, bar_y + Inches(0.32), x1 - x0, Inches(0.32))
        write(tf, lab, size=12.5, colour=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for v in (t_low, t_high):
        rect(slide, at(v) - Emu(9525), bar_y - Inches(0.22), Emu(19050),
             bar_h + Inches(0.44), INK)
        tf = txbox(slide, at(v) - Inches(0.6), bar_y + bar_h + Inches(0.26), Inches(1.2),
                   Inches(0.3))
        write(tf, f"{v} CP", size=12, bold=True, align=PP_ALIGN.CENTER)

    # A mark near either end would push its label off the page, so the label box is clamped to the
    # margins and only the leader line stays on the value.
    for m in d.get("marks", []):
        mx = at(m["at"])
        rect(slide, mx - Emu(4763), bar_y - Inches(0.80), Emu(9525), Inches(0.80), FAINT)
        lab_w = Inches(3.8)
        lx = min(max(mx - lab_w / 2, MARGIN), W - MARGIN - lab_w)
        tf = txbox(slide, lx, bar_y - Inches(1.22), lab_w, Inches(0.36))
        write(tf, m["label"], size=12, bold=True, colour=INK, align=PP_ALIGN.CENTER)
    if foot:
        tf = txbox(slide, MARGIN, bar_y + bar_h + Inches(0.70), CONTENT_W, foot_h)
        write(tf, foot, size=13, colour=INK, line=1.24)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


FIG_DIR_NAME = "02_outputs/figures_chill"

def s_problems(prs, d, page):
    """What went wrong at each step, beside what was done about it, one row per problem.

    A list of workarounds reads terribly as prose, and it is exactly what someone repeating the
    method needs most. Pairing them on a row means nobody has to hold the problem in their head
    while they read the fix, and the fixes stay short because the box will not take a paragraph.
    """
    slide = blank(prs)
    y = head(slide, d)
    items = d["items"]
    foot = d.get("foot")
    foot_h = est_h(foot, CONTENT_W, 12, line=1.2) + Inches(0.22) if foot else Inches(0)
    avail = H - y - Inches(0.58) - foot_h

    head_w = Inches(2.45)
    arrow_w = Inches(0.34)
    prob_w = int((CONTENT_W - head_w - arrow_w) * 0.42)
    sol_w = CONTENT_W - head_w - arrow_w - prob_w - Inches(0.20)
    gap = Inches(0.12)
    row_h = min(int((avail - gap * (len(items) - 1)) / len(items)),
                max(max(est_h(it["problem"], prob_w - Inches(0.28), 11.5, line=1.18),
                        est_h(it["solution"], sol_w - Inches(0.28), 11.5, line=1.18))
                    for it in items) + Inches(0.34))
    y = centred(y, avail, row_h * len(items) + gap * (len(items) - 1))

    for i, it in enumerate(items):
        ry = y + i * (row_h + gap)
        accent = it.get("colour") or ACCENTS[i % len(ACCENTS)]
        rrect(slide, MARGIN, ry, CONTENT_W, row_h, PANEL)
        rect(slide, MARGIN, ry, Inches(0.07), row_h, accent)
        tf = txbox(slide, MARGIN + Inches(0.24), ry + Inches(0.16), head_w - Inches(0.32),
                   row_h - Inches(0.3))
        write(tf, it["head"], size=12.5, bold=True, colour=accent, line=1.05)
        px = MARGIN + head_w
        tf = txbox(slide, px, ry + Inches(0.17), prob_w - Inches(0.28), row_h - Inches(0.32))
        write(tf, it["problem"], size=11.5, colour=MUTED, line=1.18)
        # The chevron is what makes the row read left to right as trouble and answer rather than
        # as two unrelated columns; it is drawn rather than typed so it cannot wrap.
        ax = px + prob_w - Inches(0.10)
        sh = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, ax, ry + int(row_h / 2) - Inches(0.11),
                                    arrow_w, Inches(0.22))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RULE
        sh.line.fill.background()
        sh.shadow.inherit = False
        sx = ax + arrow_w + Inches(0.14)
        rrect(slide, sx - Inches(0.10), ry + Inches(0.09), sol_w + Inches(0.10),
              row_h - Inches(0.18), WHITE)
        tf = txbox(slide, sx, ry + Inches(0.17), sol_w - Inches(0.18), row_h - Inches(0.32))
        write(tf, it["solution"], size=11.5, colour=INK, line=1.18)

    if foot:
        tf = txbox(slide, MARGIN, H - Inches(0.58) - foot_h, CONTENT_W, foot_h)
        write(tf, foot, size=12, colour=INK, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_funnel(prs, d, page):
    """Parallel attrition funnels, each on its own scale because the units cannot share one.

    Stations, seasons and square kilometres put on a common axis would draw a bar for the discarded
    seasons four times the width of the one for discarded stations and say nothing by it. Each
    column is therefore scaled to its own first stage, and the unit is named at the top so the
    columns are never read against each other.
    """
    slide = blank(prs)
    y = head(slide, d)
    cols = d["columns"]
    foot = d.get("foot")
    foot_h = est_h(foot, CONTENT_W, 12, line=1.2) + Inches(0.22) if foot else Inches(0)
    avail = H - y - Inches(0.58) - foot_h

    n = len(cols)
    gap = Inches(0.34)
    col_w = int((CONTENT_W - gap * (n - 1)) / n)
    rows = max(len(c["stages"]) for c in cols)
    head_h = Inches(0.58)
    stage_h = min(Inches(1.02), int((avail - head_h - Inches(0.2)) / rows))
    block_h = head_h + stage_h * rows
    y = centred(y, avail, block_h)

    for i, c in enumerate(cols):
        x = MARGIN + i * (col_w + gap)
        accent = c.get("colour") or ACCENTS[i % len(ACCENTS)]
        tf = txbox(slide, x, y, col_w, Inches(0.3))
        write(tf, c["head"], size=13, bold=True, colour=accent)
        tf = txbox(slide, x, y + Inches(0.26), col_w, Inches(0.26))
        write(tf, c["unit"], size=10, colour=FAINT)
        top = float(c["stages"][0]["value"])
        last = len(c["stages"]) - 1
        for j, st in enumerate(c["stages"]):
            sy = y + head_h + j * stage_h
            frac = max(0.06, float(st["value"]) / top) if top else 0.06
            bw = int(col_w * frac)
            solid = j in (0, last)
            # The bar is centred so the column narrows symmetrically: that taper is the whole point
            # of drawing a funnel rather than printing a table of counts.
            bx = x + int((col_w - bw) / 2)
            rect(slide, bx, sy, bw, stage_h - Inches(0.30),
                 accent if solid else PANEL, line=None if solid else RULE)
            # A stage holding a fraction of a per cent draws a bar narrower than its own figure. Put
            # the figure beside that bar in ink rather than inside it in white, where it would be
            # clipped: the whole point of the row is the number, and the bar is only its scale.
            need = Inches(0.0072 * 13 * len(st["value_text"])) + Inches(0.18)
            if bw >= need:
                tf = txbox(slide, x, sy + Inches(0.03), col_w, Inches(0.3))
                write(tf, st["value_text"], size=13, bold=True, align=PP_ALIGN.CENTER,
                      colour=WHITE if solid else INK)
            else:
                gutter = int((col_w - bw) / 2) - Inches(0.08)
                tf = txbox(slide, bx + bw + Inches(0.08), sy + Inches(0.03),
                           max(need, gutter), Inches(0.3))
                write(tf, st["value_text"], size=13, bold=True, colour=accent)
            tf = txbox(slide, x, sy + stage_h - Inches(0.27), col_w, Inches(0.26))
            write(tf, st["label"], size=9.5, colour=MUTED, align=PP_ALIGN.CENTER, line=1.08)

    if foot:
        tf = txbox(slide, MARGIN, H - Inches(0.58) - foot_h, CONTENT_W, foot_h)
        write(tf, foot, size=12, colour=INK, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


def s_decomp(prs, d, page):
    """One quantity broken into its parts, drawn to scale on a single bar.

    Written as three separate figures, a total and two subtractions of it invite the reader to
    treat them as three independent numbers. Drawn as segments of one bar they cannot be: the
    reader sees at once which part of the whole each claim is about, which is exactly what the
    difference between the two ways of quoting the recovered area turns on.
    """
    slide = blank(prs)
    y = head(slide, d)
    parts = d["parts"]
    foot = d.get("foot")
    foot_h = est_h(foot, CONTENT_W, 12, line=1.2) + Inches(0.22) if foot else Inches(0)
    avail = H - y - Inches(0.58) - foot_h

    lead = d.get("lead")
    bar_h = Inches(1.05)
    lead_h = est_h(lead, CONTENT_W, 14, line=1.2) + Inches(0.26) if lead else Inches(0)

    # Lay the labels out before drawing anything. A segment worth a tenth of the bar is narrower
    # than the words naming it, so the labels have to be allowed to overrun their segment; two of
    # them overrunning towards each other is what made the first render unreadable. Each label
    # takes the first row where it does not run into the one before it, and the last is pulled back
    # onto the page. Rows are counted here so the block can be centred at its true height.
    total = sum(float(p["value"]) for p in parts)
    lab_w = Inches(1.86)
    row_h = Inches(0.92)
    widths, places, ends = [], [], []
    x = MARGIN
    for i, p in enumerate(parts):
        w = int(CONTENT_W * float(p["value"]) / total) if total else Inches(0)
        if i == len(parts) - 1:
            w = MARGIN + CONTENT_W - x                # the last segment absorbs the rounding
        widths.append(w)
        lw = max(lab_w, w - Inches(0.12))
        lx = min(x + Inches(0.06), MARGIN + CONTENT_W - lw)
        row = 0
        while row < len(ends) and lx < ends[row]:
            row += 1
        if row == len(ends):
            ends.append(Emu(0))
        ends[row] = lx + lw + Inches(0.18)
        places.append((lx, lw, row))
        x += w
    label_h = row_h * len(ends)
    y = centred(y, avail, lead_h + Inches(0.34) + bar_h + label_h)

    if lead:
        tf = txbox(slide, MARGIN, y, CONTENT_W, lead_h)
        write(tf, lead, size=14, colour=INK, line=1.2)
        y += lead_h

    tf = txbox(slide, MARGIN, y, CONTENT_W, Inches(0.3))
    write(tf, d["total_label"], size=11, colour=FAINT)
    y += Inches(0.34)
    # Two passes on purpose. Shapes paint in insertion order, so a leader drawn for a later
    # segment would cross the label of an earlier one; drawing every bar and leader first puts all
    # of them behind every label.
    x = MARGIN
    for i, p in enumerate(parts):
        w = widths[i]
        row = places[i][2]
        accent = p.get("colour") or ACCENTS[i % len(ACCENTS)]
        rect(slide, x, y, w, bar_h, accent)
        ly = y + bar_h + Inches(0.14) + row * row_h
        rect(slide, x, y + bar_h, Emu(9525), ly - (y + bar_h), accent)
        x += w

    x = MARGIN
    for i, p in enumerate(parts):
        w = widths[i]
        lx, lw, row = places[i]
        accent = p.get("colour") or ACCENTS[i % len(ACCENTS)]
        if w >= Inches(1.0):
            tf = txbox(slide, x + Inches(0.10), y + Inches(0.16), w - Inches(0.2), Inches(0.5))
            write(tf, p["value_text"], size=15, bold=True, colour=WHITE, line=1.0)
        else:
            tf = txbox(slide, x, y + bar_h - Inches(0.44), w, Inches(0.4))
            write(tf, p["value_text"], size=12, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        ly = y + bar_h + Inches(0.14) + row * row_h
        # The label sits on a white plate so that a leader passing behind it stays hidden.
        rect(slide, lx - Inches(0.05), ly - Inches(0.04), lw, row_h - Inches(0.06), WHITE)
        tf = txbox(slide, lx, ly, lw, row_h - Inches(0.1))
        write(tf, p["label"], size=11, bold=True, colour=accent, line=1.1)
        if p.get("note"):
            para(tf, p["note"], size=10.5, colour=MUTED, line=1.16)
        x += w

    if foot:
        tf = txbox(slide, MARGIN, H - Inches(0.58) - foot_h, CONTENT_W, foot_h)
        write(tf, foot, size=12, colour=INK, line=1.2)
    footer(slide, d.get("source"), page)
    notes(slide, d.get("notes", ""))


KINDS = dict(cover=s_cover, section=s_section, figure=s_figure, figure_side=s_figure_side,
             compare=s_compare, ingredients=s_ingredients, close=s_close, gallery=s_gallery,
             stepper=s_stepper, params=s_params, twocol=s_twocol, table=s_table,
             datacard=s_datacard, bignum=s_bignum, annotated=s_annotated,
             figure_max=s_figure_max, map=s_map, figure_note=s_figure_note, flow=s_flow,
             lanes=s_lanes, cards=s_cards, timeline=s_timeline, scale=s_scale,
             problems=s_problems, funnel=s_funnel, decomp=s_decomp, closing=s_closing)


# --- the budget --------------------------------------------------------------------------------
#
# A slide is not a page. The v3 deck averaged 627 characters of on-slide text and its worst carried
# 2,010, which is a document projected rather than a slide. These two counters separate prose,
# which is read line by line and is what the budget is about, from tabular cells and axis labels,
# which are scanned. The check runs at build time so the limit cannot quietly drift back.

PROSE_KEYS = {"title", "subtitle", "lead", "body", "points", "foot", "caption", "note", "sub",
              "text", "label", "head", "value", "stat", "chip", "param", "rail", "asks",
              "problem", "solution", "value_text", "total_label", "unit", "kicker"}
NEVER = {"notes", "source", "image", "kind", "spoken", "gif", "crop", "n", "at", "notes_frac",
         "note_frac", "authors", "affil", "venue", "budget", "icon", "parts", "photo",
         "logos"}


def text_load(obj, key=None, prose=True, in_table=False):
    """Characters of on-slide text under `obj`, split into prose and tabular."""
    if isinstance(obj, str):
        if key in NEVER:
            return 0
        is_prose = key in PROSE_KEYS and not in_table
        return len(obj) if is_prose == prose else 0
    if isinstance(obj, dict):
        return sum(text_load(v, k, prose, in_table or k in ("rows", "head"))
                   for k, v in obj.items())
    if isinstance(obj, (list, tuple)):
        return sum(text_load(v, key, prose, in_table) for v in obj)
    return 0


def panels(d):
    """How many parallel panels a slide carries: cards, steps, lanes, bands, side notes.

    A lane counts its own chips too, since each one is a named thing needing a line of its own.
    """
    if isinstance(d.get("lanes"), (list, tuple)):
        return len(d["lanes"]) + sum(len(ln.get("items", ())) for ln in d["lanes"])
    for key in ("items", "steps", "bands", "notes_side", "points"):
        if isinstance(d.get(key), (list, tuple)):
            return len(d[key])
    return 1


def report_budget(deck, prose_cap, data_cap):
    """Print what each slide carries and name the ones over budget. Returns the offenders.

    The cap is not one number. A slide with a title and a sentence under it should live inside 300
    characters, but a four-card slide holding four labelled things cannot: each panel needs a name
    and a line, and squeezing them into the same total makes them cryptic rather than concise. The
    allowance therefore grows with the panel count, which keeps the pressure where it belongs, on
    slides carrying paragraphs rather than on slides carrying parts.
    """
    over = []
    total = 0
    for i, d in enumerate(deck, 1):
        # A slide may set its own allowance. A flat cap made sense while every slide argued one
        # point, but a slide whose whole job is to explain a concept to someone who has never met
        # it needs more room than a slide holding a full-page map needs, and the two cannot answer
        # to the same number. Declaring it per slide keeps the check honest: the build still fails
        # when a slide exceeds what it asked for, so the limit cannot drift by accident.
        cap = d.get("budget") or (prose_cap + 70 * max(0, panels(d) - 2))
        p = text_load(d, prose=True)
        t = text_load(d, prose=False)
        total += p + t
        if p > cap or t > data_cap:
            over.append((i, d.get("kind"), p, cap, t, d.get("title", "")[:42]))
    print(f"on-slide text: {total:,} characters over {len(deck)} slides, "
          f"mean {total // max(1, len(deck))}")
    for i, kind, p, cap, t, title in over:
        flag = f"prose {p} over {cap}" if p > cap else f"table {t} over {data_cap}"
        print(f"  slide {i:>2} [{kind}] {flag}  {title}")
    return over


def load_attrition(out):
    """Fold the two attrition tables into `out` under names a slide can quote.

    attrition_funnel_numbers.csv is shaped side/step/value and season_attrition_by_window.csv is
    one row per window, so neither fits the metric/value loader. Mapping the rows onto explicit
    names here rather than letting the content file index a table keeps the labels in the
    narrative and the arithmetic in the outputs, which is the same split every other slide uses.
    """
    steps = {}
    path = ROOT / "02_outputs" / "attrition_funnel_numbers.csv"
    if path.exists():
        with path.open(encoding="utf-8") as fh:
            for row in csv.DictReader(fh):
                steps[(row["side"], row["step"])] = float(row["value"])
    named = {
        "attr_st_served": ("stations", "Served by THREDDS"),
        "attr_st_enter": ("stations", "Enter the calculation"),
        "attr_st_archive": ("stations", "Observed archive"),
        "attr_st_api": ("stations", "Also covered by the API"),
        "attr_km2_country": ("area", "Peninsular Spain and the Balearics"),
        "attr_km2_cropland": ("area", "CORINE cropland"),
        "attr_km2_reached": ("area", "Reached by the interpolation"),
        "attr_km2_classified": ("area", "Classified into the three classes"),
    }
    for key, ident in named.items():
        if ident in steps:
            out[key] = steps[ident]

    # The seasons side is per window. The end-of-century window is the one that loses any, so it
    # is the one worth showing: everywhere else the completeness filter discards nothing at all.
    path = ROOT / "02_outputs" / "season_attrition_by_window.csv"
    if path.exists():
        with path.open(encoding="utf-8") as fh:
            rows = list(csv.DictReader(fh))
        far = next((r for r in rows if r["window"] == "far"), None)
        if far:
            out["attr_seasons_nominal"] = float(far["nominal"])
            out["attr_seasons_kept"] = float(far["kept"])
            out["attr_seasons_lost"] = float(far["lost"])
            out["attr_seasons_lost_pct"] = float(far["pct"])
        out["attr_seasons_windows_clean"] = sum(1 for r in rows if float(r["lost"]) == 0)
    return out


def load_numbers():
    """Every metric the slides quote, keyed by name, from the tables that computed them."""
    out = {}
    for name in ("talk_key_numbers.csv", "method_figure_numbers.csv", "model_spread_numbers.csv",
                 "method_chain_numbers.csv", "timeline_numbers.csv", "model_ranking_numbers.csv",
                 "cieza_numbers.csv", "v3_numbers.csv", "v3_gap_numbers.csv"):
        path = ROOT / "02_outputs" / name
        if not path.exists():
            sys.exit(f"missing {path}\n"
                     "  run scripts 27, 33, 34, 36, 37, 38, 42, 43, 45 and 47 before building the deck")
        with open(path, newline="", encoding="utf-8") as fh:
            for row in csv.DictReader(fh):
                try:
                    out[row["metric"]] = float(row["value"])
                except ValueError:
                    out[row["metric"]] = row["value"]
    return load_attrition(out)


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", default=None)
    ap.add_argument("--annex", action="store_true",
                    help="build the backup deck instead of the talk")
    ap.add_argument("--short", action="store_true",
                    help="build only the slides marked spoken=True, for the 15-minute slot")
    ap.add_argument("--v3", action="store_true",
                    help="build the methodological deck a co-author reads, not the spoken talk")
    ap.add_argument("--v4", action="store_true",
                    help="build the review deck: one idea a slide, figures at slide size")
    ap.add_argument("--v5", action="store_true",
                    help="build the master deck: introduction to conclusions, cut with --short")
    ap.add_argument("--strict", action="store_true",
                    help="fail rather than warn when a slide is over its text budget")
    a = ap.parse_args()

    if (a.v3 or a.v4) and (a.annex or a.short):
        sys.exit("--v3 and --v4 are their own decks: no annex and no spoken subset")
    # v5 is the exception: it is the master a reader goes through alone, and the fifteen-minute
    # conference talk is cut out of it with --short rather than written separately, so the two
    # cannot come to disagree the way two hand-kept lists would.
    if a.v5 and a.annex:
        sys.exit("--v5 has no annex: the backup slides live in the master itself")
    if sum((bool(a.v3), bool(a.v4), bool(a.v5))) > 1:
        sys.exit("pick one of --v3, --v4 and --v5")

    N = load_numbers()
    if a.v5:
        deck = talk_content.v5(N)
    elif a.v4:
        deck = talk_content.v4(N)
    elif a.v3:
        deck = talk_content.v3(N)
    elif a.annex:
        deck = talk_content.annex(N)
    else:
        deck = talk_content.slides(N)

    # Two decks come out of one narrative. The full one is what the coauthors review, and it has to
    # carry every check that was run; the short one is what fits a 15-minute slot. Marking the
    # spoken subset in the content file rather than keeping a second list means the two cannot
    # drift, which is the same reason the numbers come from CSVs instead of being typed.
    if a.short:
        deck = [d for d in deck if d.get("spoken")]
        if not deck:
            sys.exit("--short: no slide carries spoken=True in talk_content.py")
        # A slide may carry a second, shorter list for this cut. The close does, because the full
        # talk shows evidence the fifteen-minute one leaves out and a close should not assert what
        # the audience was never shown.
        deck = [dict(d, points=d["points_short"]) if d.get("points_short") else d for d in deck]

    if a.out is None:
        stem = ("charla_plinius_v5_15min" if (a.v5 and a.short) else
                "charla_plinius_v5" if a.v5 else
                "charla_plinius_v4" if a.v4 else
                "charla_plinius_v3" if a.v3 else
                "anexo_plinius" if a.annex else
                "charla_plinius_15min" if a.short else "charla_plinius")
        a.out = str(ROOT / "03_presentacion" / f"{stem}.pptx")

    # The budget only governs the review deck. The conference talk and its backup were finished
    # under the older layout and re-measuring them here would report a failure that is not one.
    if a.v4 or a.v5:
        over = report_budget(deck, prose_cap=420 if a.v5 else 300, data_cap=900)
        if over and a.strict:
            sys.exit(f"{len(over)} slides over budget")

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for i, d in enumerate(deck, start=1):
        kind = d["kind"]
        if kind not in KINDS:
            sys.exit(f"unknown slide kind {kind!r} on slide {i}")
        KINDS[kind](prs, d, i)

    out = Path(a.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)

    n_notes = sum(1 for s in prs.slides if s.has_notes_slide
                  and s.notes_slide.notes_text_frame.text.strip())
    print(f"{out}")
    print(f"{len(deck)} slides, {n_notes} with speaker notes, "
          f"{out.stat().st_size/1e6:.2f} MB")
    if REF_COUNT[0]:
        print(f"{REF_COUNT[0]} code references resolve to a real line of a real script")
    # The cap was 20 while the deck argued a result, 30 once it also had to explain the method, and
    # 35 once the results section showed the eleven models per scenario before their median. The v3
    # deck is read rather than delivered, so it answers to a different limit: past about 55 slides a
    # reviewer stops reading, which is the failure that matters there.
    # The v4 deck answers to a different limit again. Splitting the dense slides so each carries one
    # idea trades slide count for reading speed, and a reviewer turns pages faster than they parse
    # paragraphs, so the cap goes up while the text on each page goes down.
    cap = 80 if a.v5 else 75 if a.v4 else 55 if a.v3 else 35
    if not a.annex and len(deck) > cap:
        print(f"WARNING: {len(deck)} slides, above the cap of {cap}")


if __name__ == "__main__":
    main()
