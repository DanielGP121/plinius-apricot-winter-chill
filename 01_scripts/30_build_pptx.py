#!/usr/bin/env python3
"""Builds the projectable PowerPoint version of the deck.

Reads deck_content.json, written by 29_build_deck.R from the same narrative that produces the HTML,
so the two cannot drift apart. What differs is the treatment, not the content:

  the HTML is a document, so it carries the full prose on the page;
  this is a deck, so each slide carries the opening sentence of every paragraph and the complete
  prose goes into the speaker notes, where a presenter can actually use it.

The four interactive panels of the HTML have no equivalent here, so each becomes a static slide that
makes the same point: a plotted response curve, a drawn timeline, a two-column script map, and a bar
chart of the threshold sensitivity at the four settings that matter.

Usage:  python 30_build_pptx.py [--out ../03_presentacion/deck.pptx]
Requires: python-pptx.
"""

import json, os, sys, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import XyChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGDIR = os.path.join(ROOT, "02_outputs", "figures_chill")
PRES = os.path.join(ROOT, "03_presentacion")
OUT = sys.argv[sys.argv.index("--out") + 1] if "--out" in sys.argv else os.path.join(PRES, "deck.pptx")

# Palette taken from the viability maps themselves, so a slide and the figure on it never disagree
# about what blue, orange and red mean.
NAVY = RGBColor(0x14, 0x30, 0x4A)
BLUE = RGBColor(0x2C, 0x7B, 0xB6)
ORANGE = RGBColor(0xE6, 0x55, 0x0D)
RED = RGBColor(0xD7, 0x19, 0x1C)
INK = RGBColor(0x1D, 0x1D, 0x1F)
MUTED = RGBColor(0x6B, 0x6B, 0x70)
TINT = RGBColor(0xF2, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xC9, 0xDA, 0xE6)

HEAD, BODY = "Cambria", "Calibri"
W, H = 13.333, 7.5
M = 0.62                      # slide margin

data = json.load(open(os.path.join(PRES, "deck_content.json"), encoding="utf-8"))
DECK = data["deck"]
SWEEP = data["sweep"]
RESP = data["response"]
VIAB = data.get("viability") or []

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
BLANK = prs.slide_layouts[6]


def aslist(x):
    if x is None:
        return []
    return x if isinstance(x, list) else [x]


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY if dark else WHITE
    return s


def tbox(s, x, y, w, h, text, size=14, font=BODY, color=INK, bold=False,
         align=PP_ALIGN.LEFT, space=6, anchor=MSO_ANCHOR.TOP, italic=False):
    """One text box, one paragraph per list item. Margins zeroed so text aligns with shapes."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(aslist(text)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.name = font
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return box


def bullets(s, x, y, w, h, items, size=15):
    """Key points, one per paragraph, each led by a small square in the accent colour.

    A leading glyph rather than PowerPoint's own bullet: the built-in one inherits from a layout
    this deck does not use, and would render inconsistently across the slide kinds."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(11)
        d = p.add_run()
        d.text = "▪  "
        d.font.size = Pt(size)
        d.font.name = BODY
        d.font.color.rgb = ORANGE
        r = p.add_run()
        r.text = it
        r.font.size = Pt(size)
        r.font.name = BODY
        r.font.color.rgb = INK
    return box


def title(s, text, y=M, size=30, color=INK):
    return tbox(s, M, y, W - 2 * M, 1.0, text, size=size, font=HEAD, bold=True, color=color, space=0)


def note(s, blocks):
    txt = "\n\n".join(t for t in blocks if t)
    if txt:
        s.notes_slide.notes_text_frame.text = txt


def fig(s, name, x, y, w, h, cap=None):
    """Place a figure inside the given box, preserving aspect and centring the leftover space."""
    p = os.path.join(FIGDIR, name + ".png")
    if not os.path.exists(p):
        print("   falta figura:", name)
        return
    from PIL import Image
    iw, ih = Image.open(p).size
    ch = h - (0.26 if cap else 0)
    sc = min(w / iw, ch / ih)
    fw, fh = iw * sc, ih * sc
    s.shapes.add_picture(p, Inches(x + (w - fw) / 2), Inches(y + (ch - fh) / 2), Inches(fw), Inches(fh))
    if cap:
        tbox(s, x, y + ch, w, 0.24, cap, size=9.5, color=MUTED, align=PP_ALIGN.CENTER, space=0)


def first_sentence(t, cap=235):
    """Opening sentence of a paragraph, for the slide; the rest lives in the notes.

    Splitting on a full stop followed by a capital avoids cutting at decimals and abbreviations,
    which this text is full of (0,62 CP, Ruiz et al. 2019)."""
    parts = re.split(r"(?<=[.:])\s+(?=[A-ZÁÉÍÓÚÑ¿¡])", t.strip())
    out = parts[0].strip()
    if len(out) > cap:
        cut = out[:cap].rsplit(" ", 1)[0]
        out = cut + "…"
    return out


def card(s, x, y, w, h, fill=TINT, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# --------------------------------------------------------------------------------------------
# the four panels that replace the interactive ones
# --------------------------------------------------------------------------------------------
def panel_response(s, y):
    cd = XyChartData()
    ser = cd.add_series("Porciones por día")
    for t, c in zip(RESP["t"], RESP["cp"]):
        ser.add_data_point(t, c)
    gf = s.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
                            Inches(M), Inches(y), Inches(7.4), Inches(3.5), cd)
    ch = gf.chart
    ch.has_legend = False
    ch.has_title = False
    pl = ch.plots[0]
    pl.series[0].format.line.color.rgb = BLUE
    pl.series[0].format.line.width = Pt(2.5)
    for ax, lo, hi, t in ((ch.value_axis, 0, 0.85, "Porciones de frío por día"),
                          (ch.category_axis, -12, 22, "Temperatura media (°C)")):
        ax.minimum_scale, ax.maximum_scale = lo, hi
        ax.has_title = True
        ax.axis_title.text_frame.text = t
        for p in ax.axis_title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size, r.font.name, r.font.color.rgb = Pt(10), BODY, MUTED
        ax.tick_labels.font.size = Pt(9)
        ax.tick_labels.font.color.rgb = MUTED
        ax.has_major_gridlines = (ax is ch.value_axis)
    facts = [("−4 °C y por debajo", "cero acumulación"), ("0 °C", "11 % del óptimo"),
             ("8 °C", "óptimo, 0,79 CP/día"), ("14 °C y por encima", "cero de nuevo")]
    yy = y + 0.15
    for a, b in facts:
        card(s, 8.35, yy, 4.35, 0.72)
        tbox(s, 8.55, yy + 0.10, 4.0, 0.26, a, size=12, bold=True, color=NAVY, space=0)
        tbox(s, 8.55, yy + 0.38, 4.0, 0.26, b, size=11, color=MUTED, space=0)
        yy += 0.85


def panel_threshold(s, y):
    """Same sweep the HTML slider reads, sampled at the four settings that carry an argument."""
    # curves arrives as an object keyed by situation (a named list in R), not as an array
    cur = SWEEP["curves"]
    key = "ssp370_far"
    curve = cur[key] if isinstance(cur, dict) else cur[SWEEP["sits"].index(key)]
    thr, tot = SWEEP["thr"], SWEEP["total"]

    def F(x):
        lo, hi = 0, len(thr) - 1
        if x <= thr[0]:
            return curve[0]
        if x >= thr[hi]:
            return curve[hi]
        while hi - lo > 1:
            m = (lo + hi) // 2
            if thr[m] <= x:
                lo = m
            else:
                hi = m
        w = (x - thr[lo]) / (thr[hi] - thr[lo])
        return curve[lo] + w * (curve[hi] - curve[lo])

    cases = [("Umbrales publicados\n47,5 y 33,7", 0.0), ("Un error estándar menos\n−3,3", -3.3),
             ("Un error estándar más\n+3,3", 3.3), ("Escala de 1988\n+6,9", 6.94)]
    cd = CategoryChartData()
    cd.categories = [c[0] for c in cases]
    both = [F(47.5 + d) for _, d in cases]
    only = [F(33.7 + d) - F(47.5 + d) for _, d in cases]
    none = [tot - F(33.7 + d) for _, d in cases]
    cd.add_series("Ambas variedades", both)
    cd.add_series("Solo 'Búlida Precoz'", only)
    cd.add_series("Ninguna", none)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(M), Inches(y),
                            Inches(8.0), Inches(3.6), cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(10)
    for ser, col in zip(ch.plots[0].series, (BLUE, RGBColor(0xFD, 0xAE, 0x61), RED)):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = col
    ch.value_axis.tick_labels.font.size = Pt(9)
    ch.value_axis.tick_labels.font.color.rgb = MUTED
    ch.value_axis.has_title = True
    ch.value_axis.axis_title.text_frame.text = "km² de suelo cultivable"
    for p in ch.value_axis.axis_title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size, r.font.name, r.font.color.rgb = Pt(10), BODY, MUTED
    ch.category_axis.tick_labels.font.size = Pt(9)
    ch.category_axis.tick_labels.font.color.rgb = MUTED

    card(s, 8.85, y + 0.1, 3.85, 3.1, fill=TINT)
    tbox(s, 9.1, y + 0.32, 3.4, 0.3, "La banda del mutante", size=13, bold=True, color=NAVY, space=0)
    rows = [(f"{only[0]:,.0f} km²".replace(",", "."), "con los umbrales publicados"),
            (f"{only[1]:,.0f} km²".replace(",", "."), "un error estándar por debajo"),
            (f"{only[2]:,.0f} km²".replace(",", "."), "un error estándar por encima"),
            (f"{only[3]:,.0f} km²".replace(",", "."), "en la escala de 1988")]
    yy = y + 0.75
    for a, b in rows:
        tbox(s, 9.1, yy, 3.4, 0.28, a, size=15, bold=True, color=ORANGE, space=0)
        tbox(s, 9.1, yy + 0.27, 3.4, 0.24, b, size=10, color=MUTED, space=0)
        yy += 0.6


def panel_timeline(s, y):
    x0, wtot, y0 = M + 2.55, W - 2 * M - 2.75, y + 0.2
    Y0, Y1 = 1975, 2100
    rows = [("Observado, archivo PNACC", 1975, 2020, BLUE, "3044 estaciones"),
            ("Observado, API AEMET", 1995, 2025, RGBColor(0x5A, 0xA9, 0xD6), "666 estaciones"),
            ("CMIP6 histórico", 1975, 2014, RGBColor(0x8A, 0x8A, 0x90), "11 modelos"),
            ("CMIP6 escenarios SSP", 2015, 2100, RGBColor(0xA8, 0xA8, 0xAE), "3 × 11"),
            ("Línea base 1995-2020", 1995, 2020, ORANGE, "26 temporadas"),
            ("Próximo plazo 2021-2040", 2021, 2040, ORANGE, "20 temporadas"),
            ("Medio siglo 2041-2070", 2041, 2070, ORANGE, "30 temporadas"),
            ("Fin de siglo 2071-2100", 2071, 2100, ORANGE, "30 temporadas")]
    xs = lambda yr: x0 + (yr - Y0) / (Y1 - Y0) * wtot
    for i, (n, a, b, c, note_) in enumerate(rows):
        yy = y0 + i * 0.44
        tbox(s, M, yy, 2.45, 0.32, n, size=11, color=INK, align=PP_ALIGN.RIGHT, space=0)
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(xs(a)), Inches(yy),
                                Inches(max(xs(b) - xs(a), 0.12)), Inches(0.26))
        sh.adjustments[0] = 0.3
        sh.fill.solid()
        sh.fill.fore_color.rgb = c
        sh.line.fill.background()
        sh.shadow.inherit = False
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = note_
        r.font.size, r.font.name, r.font.color.rgb, r.font.bold = Pt(9), BODY, WHITE, True
    ybot = y0 + len(rows) * 0.44
    for yr in (1980, 2000, 2020, 2040, 2060, 2080, 2100):
        tbox(s, xs(yr) - 0.3, ybot + 0.02, 0.6, 0.24, str(yr), size=9, color=MUTED,
             align=PP_ALIGN.CENTER, space=0)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(xs(2014.5)), Inches(y0 - 0.12),
                            Inches(0.018), Inches(ybot - y0 + 0.12))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RED
    ln.line.fill.background()
    ln.shadow.inherit = False
    tbox(s, xs(2014.5) - 1.1, y0 - 0.42, 2.2, 0.26, "costura histórico / SSP", size=9.5,
         color=RED, align=PP_ALIGN.CENTER, space=0)


def panel_workflow(s, y):
    L = [("14_ladon_download_thredds.sh", "88 NetCDF desde THREDDS, ~15 GB"),
         ("15_chill_national_parallel.R", "motor nacional de frío, checkpoints y empalme"),
         ("21_aemet_observed_download.py", "observado reciente por API")]
    R = [("22_merge_chill_tables.R", "tabla canónica de 462.808 filas"),
         ("19 y 28_cropland", "IDW + CORINE, mapas y barrido de umbrales"),
         ("23 a 25_observado", "frío de la API, contraste y empalme"),
         ("26 y 27_registro", "1976-2025 y comprobación en Cieza"),
         ("29 y 30_build_deck", "este documento y su versión HTML")]
    for col, (items, head, x) in enumerate(((L, "HPC Ladon", M), (R, "Local", 6.95))):
        tbox(s, x, y, 5.7, 0.3, head.upper(), size=11, bold=True, color=ORANGE, space=0)
        # 0.80" pitch, not 0.90": the local column holds five cards and the panel starts low enough
        # that the looser spacing pushed the last one off the bottom of the slide
        yy = y + 0.40
        for a, b in items:
            card(s, x, yy, 5.75, 0.70)
            tbox(s, x + 0.22, yy + 0.09, 5.3, 0.24, a, size=11.5, bold=True, color=NAVY, space=0)
            tbox(s, x + 0.22, yy + 0.38, 5.3, 0.24, b, size=10.5, color=MUTED, space=0)
            yy += 0.80


PANELS = {"response": panel_response, "threshold": panel_threshold,
          "timeline": panel_timeline, "workflow": panel_workflow}

# --------------------------------------------------------------------------------------------
# build
# --------------------------------------------------------------------------------------------
n_fig = 0
for b in DECK:
    kind = b["kind"]

    if kind == "cover":
        s = slide(dark=True)
        tbox(s, M, 1.75, W - 2 * M - 1.2, 2.2, b["title"], size=38, font=HEAD, bold=True, color=WHITE, space=10)
        tbox(s, M, 4.15, W - 2 * M - 1.2, 0.7, b["subtitle"], size=17, color=PALE, space=0)
        tbox(s, M, 5.5, W - 2 * M, 1.4, aslist(b["meta"]), size=11.5, color=PALE, space=3)

    elif kind == "section":
        s = slide(dark=True)
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(M), Inches(2.5), Inches(1.15), Inches(1.15))
        c.fill.solid()
        c.fill.fore_color.rgb = ORANGE
        c.line.fill.background()
        c.shadow.inherit = False
        tf = c.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = str(b["n"])
        r.font.size, r.font.name, r.font.bold, r.font.color.rgb = Pt(30), HEAD, True, WHITE
        # 1.35" of title, not 1.0": several part names wrap to two lines at 34pt
        tbox(s, M + 1.6, 2.35, W - M - 2.2, 1.35, b["title"], size=34, font=HEAD, bold=True, color=WHITE, space=0)
        tbox(s, M + 1.6, 3.85, W - M - 3.0, 1.6, b["lead"], size=15, color=PALE, space=0)

    elif kind == "gadget":
        s = slide()
        title(s, b["title"])
        pts = [first_sentence(t, cap=175) for t in aslist(b["body"])][:2]
        bullets(s, M, 1.5, W - 2 * M, 1.3, pts, size=13)
        PANELS[b["id"]](s, 2.9)
        note(s, aslist(b["body"]) + [("NOTA. " + b["note"]) if b.get("note") else ""])

    elif kind == "gallery":
        figs = aslist(b["figs"])
        per = 6
        for k in range(0, len(figs), per):
            s = slide()
            part = f"  ({k // per + 1}/{(len(figs) + per - 1) // per})" if len(figs) > per else ""
            title(s, b["title"] + part, size=26)
            if k == 0:
                tbox(s, M, 1.32, W - 2 * M, 0.4, aslist(b["body"])[0], size=12, color=MUTED, space=0)
            gw, gh = (W - 2 * M - 0.5) / 3, 2.25
            for j, f in enumerate(figs[k:k + per]):
                fig(s, f, M + (j % 3) * (gw + 0.25), 1.85 + (j // 3) * (gh + 0.25), gw, gh,
                    cap=f.replace("_", " "))
                n_fig += 1

    else:  # slide
        s = slide()
        title(s, b["title"])
        figs = aslist(b.get("figs"))
        caps = aslist(b.get("figcap"))
        pts = [first_sentence(t) for t in aslist(b["body"])]
        if b.get("table") == "viability" and VIAB:
            bullets(s, M, 1.55, W - 2 * M, 1.5, pts[:2], size=14)
            rowsn = min(len(VIAB), 11) + 1
            tb = s.shapes.add_table(rowsn, 4, Inches(M), Inches(3.2), Inches(W - 2 * M), Inches(3.75)).table
            hdr = ["Situación", "Ambas", "Solo 'Búlida Precoz'", "Ninguna"]
            for j, htxt in enumerate(hdr):
                cell = tb.cell(0, j)
                cell.text = htxt
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size, r.font.bold, r.font.name = Pt(11), True, BODY
            for i, v in enumerate(VIAB[:rowsn - 1]):
                vals = [v["label"], f"{v['pct_both']:.1f}", f"{v['pct_only_precoz']:.1f}", f"{v['pct_none']:.1f}"]
                for j, val in enumerate(vals):
                    cell = tb.cell(i + 1, j)
                    cell.text = val
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.size, r.font.name = Pt(10), BODY
                            if j == 2:
                                r.font.bold, r.font.color.rgb = True, ORANGE
        elif not figs:
            half = (len(pts) + 1) // 2
            bullets(s, M, 1.6, (W - 2 * M - 0.6) / 2, 4.6, pts[:half])
            bullets(s, M + (W - 2 * M - 0.6) / 2 + 0.6, 1.6, (W - 2 * M - 0.6) / 2, 4.6, pts[half:])
        else:
            bullets(s, M, 1.6, 5.55, 4.6, pts, size=14)
            fx, fw = 6.55, W - M - 6.55
            if len(figs) == 1:
                fig(s, figs[0], fx, 1.55, fw, 4.75, cap=caps[0] if caps else None)
            else:
                fh = 4.75 / len(figs) - 0.15
                for j, f in enumerate(figs[:3]):
                    fig(s, f, fx, 1.55 + j * (fh + 0.15), fw, fh, cap=caps[j] if len(caps) > j else None)
            n_fig += len(figs)
        note(s, aslist(b["body"]) + [("NOTA. " + b["note"]) if b.get("note") else ""])

os.makedirs(PRES, exist_ok=True)
prs.save(OUT)
print(f"escrito {OUT}")
print(f"  {len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas, {n_fig} figuras, "
      f"{os.path.getsize(OUT)/1048576:.1f} MB")
